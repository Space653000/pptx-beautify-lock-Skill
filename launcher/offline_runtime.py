from __future__ import annotations

from pathlib import Path
import hashlib
import os
import uuid

from pptx import Presentation

import pptx_offline_engine as bundled_engine
from update_manager import check_for_engine_update, load_engine_module

LAUNCHER_VERSION = "0.7.1"
BUNDLED_ENGINE_VERSION = "0.7.1"


def _verify_pptx(path: Path, expected_slides: int) -> tuple[int, str]:
    if not path.is_file():
        raise RuntimeError(f"OUTPUT_MISSING: {path}")
    size = path.stat().st_size
    if size <= 0:
        raise RuntimeError(f"OUTPUT_EMPTY: {path}")
    reopened = Presentation(str(path))
    if len(reopened.slides) != expected_slides:
        raise RuntimeError(
            f"OUTPUT_REOPEN_SLIDE_COUNT_MISMATCH: expected={expected_slides} actual={len(reopened.slides)}"
        )
    digest = hashlib.sha256(path.read_bytes()).hexdigest().upper()
    return size, digest


def beautify_to_final(
    source: str | Path,
    output: str | Path,
    style: str,
    log=lambda _text: None,
    *,
    check_updates: bool = True,
):
    src = Path(source)
    out = Path(output)
    if not src.is_file():
        raise ValueError(f"來源不存在：{src}")
    if src.suffix.lower() != ".pptx" or out.suffix.lower() != ".pptx":
        raise ValueError("輸入與輸出都必須是 .pptx")
    if src.resolve() == out.resolve():
        raise ValueError("輸出不可覆寫來源 PPTX")

    expected_slides = len(Presentation(str(src)).slides)
    out.parent.mkdir(parents=True, exist_ok=True)

    selection = None
    engine = bundled_engine
    if check_updates:
        selection = check_for_engine_update(
            local_version=BUNDLED_ENGINE_VERSION,
            launcher_version=LAUNCHER_VERSION,
            log=log,
        )
        engine = load_engine_module(selection, bundled_engine, log)
    else:
        log("UPDATE_CHECK=selftest_disabled")
        log(f"EFFECTIVE_ENGINE_VERSION={BUNDLED_ENGINE_VERSION}")

    candidate = out.with_name(f".{out.stem}.{uuid.uuid4().hex}.candidate.pptx")

    def engine_log(message: str):
        # The underlying engine validates the candidate, but only the wrapper may
        # claim final delivery success after atomic promotion and reopen checks.
        if message.strip() == "OFFLINE_BEAUTIFY_PASS=true":
            return
        log(message)

    try:
        report = engine.beautify_pptx(src, candidate, style, engine_log)
        candidate_size, _ = _verify_pptx(candidate, expected_slides)
        log("CANDIDATE_REOPEN_PASS=true")
        log(f"CANDIDATE_BYTES={candidate_size}")

        os.replace(candidate, out)

        final_size, final_sha256 = _verify_pptx(out, expected_slides)
        log("FINAL_OUTPUT_EXISTS=true")
        log("FINAL_OUTPUT_REOPEN_PASS=true")
        log(f"FINAL_OUTPUT_BYTES={final_size}")
        log(f"FINAL_OUTPUT_SHA256={final_sha256}")
        log(f"FINAL_OUTPUT_PATH={out.resolve()}")
        log("OFFLINE_BEAUTIFY_PASS=true")
        return report
    except Exception:
        if candidate.exists():
            candidate.unlink(missing_ok=True)
        raise
