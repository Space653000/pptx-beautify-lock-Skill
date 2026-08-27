#!/usr/bin/env python3
"""PyInstaller entry point for the offline-first PPTX beautifier."""
from __future__ import annotations

from pathlib import Path
import shutil
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches

import pptx_beautify_gui as core
from offline_runtime import beautify_to_final


def _portable_self_test() -> int:
    try:
        if not core.BEAUTIFY_OFFLINE or core.CLOUD_AI_ENABLED or core.NETWORK_REQUIRED:
            return 10
        if not core.OPTIONAL_UPDATE_CHECK:
            return 11
        if tuple(core.PRODUCT_FEATURES) != ("input_pptx", "output_pptx", "style", "beautify"):
            return 12

        tmp = Path(tempfile.mkdtemp(prefix="pptx-offline-selftest-"))
        try:
            src = tmp / "source.pptx"
            out = tmp / "out.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Offline Beautifier Self Test"
            table = slide.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(8), Inches(2)).table
            for r in range(3):
                for c in range(3):
                    table.cell(r, c).text = f"R{r}C{c}"
            prs.save(src)

            logs: list[str] = []
            beautify_to_final(
                src,
                out,
                "專業技術（Technical Clean）",
                logs.append,
                check_updates=False,
            )
            if not out.is_file() or out.stat().st_size < 1000:
                return 13
            if "FINAL_OUTPUT_EXISTS=true" not in logs:
                return 14
            if "FINAL_OUTPUT_REOPEN_PASS=true" not in logs:
                return 15
            if logs.count("OFFLINE_BEAUTIFY_PASS=true") != 1:
                return 16
            check = Presentation(out)
            if check.slides[0].shapes.title.text != "Offline Beautifier Self Test":
                return 17
        finally:
            shutil.rmtree(tmp, ignore_errors=True)
        return 0
    except Exception:
        return 18


if __name__ == "__main__":
    if "--portable-self-test" in sys.argv:
        raise SystemExit(_portable_self_test())
    core.App().mainloop()
