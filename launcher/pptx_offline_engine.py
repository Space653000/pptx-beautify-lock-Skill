from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
from io import BytesIO

from PIL import Image, ImageStat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt

try:
    from pptx_content_lock import build_manifest, diff as content_diff
except Exception:
    build_manifest = None
    content_diff = None

CJK_RE = re.compile(r"[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]")
BRAND_WORDS = (
    "PEGATRON", "MICROSOFT", "SURFACE", "MEC",
    "MECHANICAL ENGINEERING CENTER", "SUPPORT, INTEGRATION", "COPYRIGHT", "©",
)
GENERIC_TEMPLATE_TEXT = (
    "EVENT NAME OR PRESENTATION TITLE", "SPEAKER NAME OR SUBTITLE",
    "PRESENTATION TITLE", "CLICK TO ADD TITLE", "CLICK TO ADD SUBTITLE",
    "CLICK TO ADD TEXT", "CLICK TO ADD CONTENT",
)

STYLE_PRESETS = {
    "自動（忠於原稿 / Source-faithful）": {
        "accent": "source", "title": RGBColor(17, 24, 39),
        "body": RGBColor(31, 41, 55), "surface": RGBColor(248, 250, 252),
        "grid": RGBColor(203, 213, 225),
    },
    "專業技術（Technical Clean）": {
        "accent": RGBColor(8, 145, 178), "title": RGBColor(15, 23, 42),
        "body": RGBColor(30, 41, 59), "surface": RGBColor(248, 250, 252),
        "grid": RGBColor(148, 163, 184),
    },
    "商務簡潔（Executive Minimal）": {
        "accent": RGBColor(37, 99, 235), "title": RGBColor(17, 24, 39),
        "body": RGBColor(55, 65, 81), "surface": RGBColor(249, 250, 251),
        "grid": RGBColor(209, 213, 219),
    },
    "現代極簡（Modern Minimal）": {
        "accent": RGBColor(17, 24, 39), "title": RGBColor(17, 24, 39),
        "body": RGBColor(75, 85, 99), "surface": RGBColor(250, 250, 250),
        "grid": RGBColor(229, 231, 235),
    },
    "高階科技簡報（Premium Tech, preserve source palette）": {
        "accent": RGBColor(6, 182, 212), "title": RGBColor(15, 23, 42),
        "body": RGBColor(30, 41, 59), "surface": RGBColor(241, 245, 249),
        "grid": RGBColor(148, 163, 184),
    },
}


@dataclass
class RunReport:
    slide_count: int
    removed_empty_placeholders: int = 0
    suppressed_template_artifacts: int = 0
    tables_styled: int = 0
    data_slides_structured: int = 0
    text_runs_normalized: int = 0
    warnings: list[str] | None = None

    def __post_init__(self):
        if self.warnings is None:
            self.warnings = []


def _has_text(shape) -> bool:
    return bool(getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())


def _is_brand(shape) -> bool:
    text = (getattr(shape, "text", "") or "").upper()
    return any(token in text for token in BRAND_WORDS)


def _is_generic_template(shape) -> bool:
    text = (getattr(shape, "text", "") or "").strip().upper()
    return any(token in text for token in GENERIC_TEMPLATE_TEXT)


def _contains_cjk(text: str) -> bool:
    return bool(CJK_RE.search(text or ""))


def _safe_font(text: str) -> str:
    return "Microsoft JhengHei" if _contains_cjk(text) else "Aptos"


def _slide_is_dark(slide, prs: Presentation) -> bool:
    sw, sh = prs.slide_width, prs.slide_height
    shape_sets = [slide.shapes]
    try:
        shape_sets.append(slide.slide_layout.shapes)
        shape_sets.append(slide.slide_layout.slide_master.shapes)
    except Exception:
        pass
    for shapes in shape_sets:
        for shape in shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if shape.width * shape.height < sw * sh * 0.62:
                continue
            try:
                im = Image.open(BytesIO(shape.image.blob)).convert("RGB")
                im.thumbnail((96, 96))
                r, g, b = ImageStat.Stat(im).mean[:3]
                lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
                return lum < 120
            except Exception:
                pass
    try:
        fill = slide.background.fill
        if fill.type and fill.fore_color.rgb:
            r, g, b = tuple(fill.fore_color.rgb)
            return 0.2126 * r + 0.7152 * g + 0.0722 * b < 120
    except Exception:
        pass
    return False


def _slide_cfg(base_cfg, dark: bool):
    cfg = dict(base_cfg)
    if dark:
        cfg["title"] = RGBColor(248, 250, 252)
        cfg["body"] = RGBColor(226, 232, 240)
        cfg["surface"] = RGBColor(30, 41, 59)
        cfg["grid"] = RGBColor(71, 85, 105)
    return cfg


def _font_sizes(shape) -> list[float]:
    vals = []
    if not getattr(shape, "has_text_frame", False):
        return vals
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                vals.append(r.font.size.pt)
    return vals


def _avg_font_size(shape) -> float:
    vals = _font_sizes(shape)
    return sum(vals) / len(vals) if vals else 0.0


def _source_accent(prs: Presentation) -> RGBColor:
    counts: dict[tuple[int, int, int], int] = {}
    for slide in list(prs.slides)[: min(6, len(prs.slides))]:
        for shape in slide.shapes:
            try:
                fill = shape.fill
                if fill.type and fill.fore_color.type is not None and fill.fore_color.rgb:
                    rgb = tuple(fill.fore_color.rgb)
                    if max(rgb) - min(rgb) > 55 and 35 < sum(rgb) / 3 < 220:
                        counts[rgb] = counts.get(rgb, 0) + 1
            except Exception:
                pass
    if counts:
        rgb = max(counts.items(), key=lambda kv: kv[1])[0]
        return RGBColor(*rgb)
    return RGBColor(6, 169, 201)


def _style_cfg(prs: Presentation, style: str):
    base = STYLE_PRESETS.get(style, STYLE_PRESETS["自動（忠於原稿 / Source-faithful）"]).copy()
    if base["accent"] == "source" or "preserve source palette" in style:
        base["accent"] = _source_accent(prs)
    return base


def _set_cell_border(cell, color: RGBColor, width: str = "6350"):
    tcPr = cell._tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        old = tcPr.find(edge, tcPr.nsmap)
        if old is not None:
            tcPr.remove(old)
        xml = (
            f'<{edge} w="{width}" {nsdecls("a")}>'
            f'<a:solidFill><a:srgbClr val="{str(color)}"/></a:solidFill>'
            f'<a:prstDash val="solid"/></{edge}>'
        )
        tcPr.append(parse_xml(xml))


def _normalize_text(shape, *, role: str, cfg, report: RunReport):
    if not getattr(shape, "has_text_frame", False) or not shape.text_frame.text.strip():
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    for p in tf.paragraphs:
        for run in p.runs:
            if not run.text:
                continue
            run.font.name = _safe_font(run.text)
            report.text_runs_normalized += 1
            if role == "title":
                if not run.font.size or run.font.size.pt < 22:
                    run.font.size = Pt(26)
                run.font.bold = True
                run.font.color.rgb = cfg["title"]
            elif role == "section":
                if not run.font.size or run.font.size.pt < 26:
                    run.font.size = Pt(30)
                run.font.bold = True
                run.font.color.rgb = cfg["title"]
            elif role == "body":
                if run.font.size and run.font.size.pt < 9:
                    run.font.size = Pt(9)
                run.font.color.rgb = cfg["body"]


def _role_for_text_shape(shape, slide_w: int, slide_h: int) -> str:
    if not _has_text(shape):
        return "none"
    text = shape.text.strip()
    if _is_brand(shape):
        return "brand"
    top_ratio = shape.top / slide_h
    width_ratio = shape.width / slide_w
    avg = _avg_font_size(shape)
    try:
        if shape.is_placeholder and shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE
        ):
            return "title"
    except Exception:
        pass
    if top_ratio < 0.13 and (avg >= 18 or width_ratio > 0.35):
        return "title"
    if len(text) < 80 and avg >= 24 and 0.15 < top_ratio < 0.65:
        return "section"
    return "body"


def _remove_empty_placeholders(slide) -> int:
    removed = 0
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False) or shape.text_frame.text.strip():
            continue
        try:
            ptype = shape.placeholder_format.type
        except Exception:
            continue
        if ptype in (
            PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT,
        ):
            sp = shape._element
            sp.getparent().remove(sp)
            removed += 1
    return removed


def _suppress_generic_template_text(slide, prs: Presentation) -> int:
    count = 0
    for shape in slide.shapes:
        if _has_text(shape) and _is_generic_template(shape):
            shape.left = prs.slide_width + Inches(0.5)
            shape.top = prs.slide_height + Inches(0.5)
            shape.width = Inches(0.1)
            shape.height = Inches(0.1)
            count += 1
    return count


def _cover_structure(slide, prs: Presentation, cfg, report: RunReport) -> bool:
    sw, sh = prs.slide_width, prs.slide_height
    texts = [s for s in slide.shapes if _has_text(s)]
    content = [s for s in texts if not _is_brand(s) and not _is_generic_template(s)]
    if not content:
        return False
    candidates = [s for s in content if len(s.text.strip()) >= 10 and _avg_font_size(s) >= 14]
    if not candidates:
        candidates = [s for s in content if len(s.text.strip()) >= 10]
    if not candidates:
        return False
    title = max(candidates, key=lambda s: (_avg_font_size(s), len(s.text), s.width))
    others = [s for s in content if s is not title]
    meta = [s for s in others if s.top < sh * 0.76 and len(s.text.strip()) < 80]

    brand_texts = [s for s in texts if _is_brand(s) and s.top < sh * 0.65]
    brand_bottom = max((s.top + s.height for s in brand_texts), default=int(sh * 0.18))
    dark = _slide_is_dark(slide, prs)
    if dark:
        y = int(sh * 0.56)
    elif brand_bottom > sh * 0.30:
        y = max(int(sh * 0.54), brand_bottom + int(sh * 0.025))
    else:
        y = int(sh * 0.37)
    y = min(y, int(sh * 0.64))

    title.left = int(sw * 0.075)
    title.top = y
    title.width = int(sw * 0.70)
    title.height = int(sh * 0.15)
    title.text_frame.word_wrap = True
    for p in title.text_frame.paragraphs:
        for run in p.runs:
            if run.text:
                run.font.name = _safe_font(run.text)
                run.font.size = Pt(26 if dark else 28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255) if dark else cfg["title"]

    cursor = title.top + title.height + int(sh * 0.015)
    for s in sorted(meta, key=lambda x: (x.top, x.left))[:2]:
        s.left = title.left
        s.top = cursor
        s.width = int(sw * 0.48)
        s.height = int(sh * 0.045)
        for p in s.text_frame.paragraphs:
            for run in p.runs:
                if run.text:
                    run.font.name = _safe_font(run.text)
                    run.font.size = Pt(12)
                    run.font.bold = False
                    run.font.color.rgb = cfg["accent"] if dark else cfg["body"]
        cursor += int(sh * 0.052)
    return True


def _style_table(shape, cfg, report: RunReport):
    tbl = shape.table
    rows, cols = len(tbl.rows), len(tbl.columns)
    if rows == 0 or cols == 0:
        return
    accent, grid, surface = cfg["accent"], cfg["grid"], cfg["surface"]
    section_rows = set()
    for r in range(rows):
        vals = [tbl.cell(r, c).text.strip() for c in range(cols)]
        if vals and vals[0] and all(not x for x in vals[1:]):
            section_rows.add(r)
    if 0 not in section_rows:
        section_rows.add(0)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.015)
            cell.margin_bottom = Inches(0.015)
            if r in section_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = accent
                fg, bold = RGBColor(255, 255, 255), True
            elif r == 1 or (r > 0 and r - 1 in section_rows):
                cell.fill.solid(); cell.fill.fore_color.rgb = surface
                fg, bold = cfg["title"], True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                fg, bold = cfg["body"], False
            _set_cell_border(cell, grid)
            tf = cell.text_frame
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    if run.text:
                        run.font.name = _safe_font(run.text)
                        run.font.size = Pt(8.5 if c == 0 else 9)
                        run.font.bold = bold
                        run.font.color.rgb = fg
                        report.text_runs_normalized += 1
    report.tables_styled += 1


def _data_slide_structure(slide, prs: Presentation, cfg, report: RunReport) -> bool:
    slide_w, slide_h = prs.slide_width, prs.slide_height
    tables = [s for s in slide.shapes if getattr(s, "has_table", False)]
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    substantial = [p for p in pics if p.width * p.height > slide_w * slide_h * 0.055]
    if len(tables) != 1 or len(substantial) != 2:
        return False

    table = tables[0]
    text_shapes = [s for s in slide.shapes if _has_text(s) and not _is_brand(s)]
    title_candidates = [s for s in text_shapes if s.top < slide_h * 0.18 and s.width > slide_w * 0.28]
    title = min(title_candidates, key=lambda s: (s.top, -s.width)) if title_candidates else None
    summaries = [
        s for s in text_shapes
        if s is not title and s.left < slide_w * 0.40 and s.top < slide_h * 0.42
        and len(s.text.strip()) >= 12
    ]
    chart_labels = [
        s for s in text_shapes
        if s is not title and len(s.text.strip()) < 36
        and re.match(r"^[LR](?:[- _]|$)", s.text.strip(), re.I)
    ]

    left_margin = int(slide_w * 0.04)
    right_margin = int(slide_w * 0.04)
    gutter = int(slide_w * 0.035)
    col_w = int((slide_w - left_margin - right_margin - gutter) / 2)

    if title is not None:
        title.left = left_margin
        title.top = int(slide_h * 0.045)
        title.width = int(slide_w * 0.68)
        title.height = max(title.height, int(slide_h * 0.09))
    if summaries:
        summary = max(summaries, key=lambda s: len(s.text))
        summary.left = left_margin
        summary.top = int(slide_h * 0.18)
        summary.width = int(slide_w * 0.35)
        summary.height = int(slide_h * 0.19)
    table.left = int(slide_w * 0.40)
    table.top = int(slide_h * 0.145)
    table.width = int(slide_w * 0.55)
    table.height = int(slide_h * 0.27)

    substantial.sort(key=lambda p: p.left)
    for i, pic in enumerate(substantial):
        pic.left = left_margin + i * (col_w + gutter)
        pic.top = int(slide_h * 0.54)
        pic.width = col_w
        pic.height = int(slide_h * 0.29)
    if len(chart_labels) >= 2:
        chart_labels.sort(key=lambda s: s.left)
        for i, lab in enumerate(chart_labels[:2]):
            lab.left = left_margin + i * (col_w + gutter)
            lab.top = int(slide_h * 0.49)
            lab.width = col_w
            lab.height = int(slide_h * 0.04)

    report.data_slides_structured += 1
    return True


def _clamp_geometry(slide, prs: Presentation):
    sw, sh = prs.slide_width, prs.slide_height
    for shape in slide.shapes:
        if _is_brand(shape) or _is_generic_template(shape):
            continue
        if shape.left < -int(sw * 0.02):
            shape.left = 0
        if shape.top < -int(sh * 0.02):
            shape.top = 0
        if shape.left + shape.width > sw * 1.02 and shape.left < sw:
            shape.width = max(Inches(0.1), sw - shape.left)
        if shape.top + shape.height > sh * 1.02 and shape.top < sh:
            shape.height = max(Inches(0.1), sh - shape.top)


def beautify_pptx(source: str | Path, output: str | Path, style: str, log=lambda x: None) -> RunReport:
    src, out = Path(source), Path(output)
    if not src.is_file() or src.suffix.lower() != ".pptx":
        raise ValueError("來源必須是有效的 .pptx")
    if out.suffix.lower() != ".pptx":
        raise ValueError("輸出必須是 .pptx")
    if src.resolve() == out.resolve():
        raise ValueError("輸出不可覆寫來源 PPTX")

    if build_manifest is None or content_diff is None:
        raise RuntimeError("CONTENT_LOCK_HELPER_MISSING: 離線版禁止在無 Content Lock 驗證下執行")
    before = build_manifest(str(src))
    prs = Presentation(str(src))
    cfg = _style_cfg(prs, style)
    report = RunReport(slide_count=len(prs.slides))

    log("OFFLINE_ENGINE=true")
    log(f"slides={len(prs.slides)}")
    log(f"style={style}")

    for idx, slide in enumerate(prs.slides, 1):
        report.removed_empty_placeholders += _remove_empty_placeholders(slide)
        report.suppressed_template_artifacts += _suppress_generic_template_text(slide, prs)
        scfg = _slide_cfg(cfg, _slide_is_dark(slide, prs))
        cover_structured = _cover_structure(slide, prs, scfg, report) if idx == 1 else False
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                _style_table(shape, scfg, report)
        structured = _data_slide_structure(slide, prs, scfg, report) if idx != 1 else False
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                role = _role_for_text_shape(shape, prs.slide_width, prs.slide_height)
                if role not in ("brand", "none"):
                    if idx == 1 and (cover_structured or _is_generic_template(shape)):
                        continue
                    _normalize_text(shape, role=role, cfg=scfg, report=report)
        _clamp_geometry(slide, prs)
        log(f"slide={idx}/{len(prs.slides)} structured={'true' if (structured or cover_structured) else 'false'}")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    after = build_manifest(str(out))
    diffs = content_diff(before, after)
    if diffs:
        try:
            out.unlink()
        except OSError:
            pass
        preview = "\n".join(diffs[:8])
        raise RuntimeError(f"CONTENT_LOCK_FAIL: {len(diffs)} differences\n{preview}")

    log("CONTENT_LOCK_PASS=true")
    log(f"empty_placeholders_removed={report.removed_empty_placeholders}")
    log(f"template_artifacts_suppressed={report.suppressed_template_artifacts}")
    log(f"tables_styled={report.tables_styled}")
    log(f"data_slides_structured={report.data_slides_structured}")
    log("OFFLINE_BEAUTIFY_PASS=true")
    return report
