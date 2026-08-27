from __future__ import annotations

from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

import pptx_offline_engine as base_engine

ENGINE_VERSION = "0.7.2"
STYLE_PRESETS = base_engine.STYLE_PRESETS
RunReport = base_engine.RunReport
SOURCE_FAITHFUL_STYLE = "自動（忠於原稿 / Source-faithful）"


def _norm_text(value: str) -> str:
    return "\n".join(line.rstrip() for line in (value or "").replace("\r\n", "\n").split("\n")).strip()


def _text_shapes(slide):
    return [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and _norm_text(shape.text_frame.text)
    ]


def _table_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_table", False)]


def _shape_distance(a, b) -> int:
    return (
        abs(int(a.left) - int(b.left))
        + abs(int(a.top) - int(b.top))
        + abs(int(a.width) - int(b.width))
        + abs(int(a.height) - int(b.height))
    )


def _match_text_shapes(source_slide, output_slide):
    buckets = defaultdict(list)
    for shape in _text_shapes(output_slide):
        buckets[_norm_text(shape.text_frame.text)].append(shape)

    pairs = []
    for source_shape in _text_shapes(source_slide):
        key = _norm_text(source_shape.text_frame.text)
        candidates = buckets.get(key, [])
        if not candidates:
            continue
        best = min(candidates, key=lambda item: _shape_distance(source_shape, item))
        candidates.remove(best)
        pairs.append((source_shape, best))
    return pairs


def _table_key(shape) -> tuple:
    table = shape.table
    return tuple(
        tuple(_norm_text(table.cell(r, c).text) for c in range(len(table.columns)))
        for r in range(len(table.rows))
    )


def _match_tables(source_slide, output_slide):
    buckets = defaultdict(list)
    for shape in _table_shapes(output_slide):
        buckets[_table_key(shape)].append(shape)

    pairs = []
    for source_shape in _table_shapes(source_slide):
        key = _table_key(source_shape)
        candidates = buckets.get(key, [])
        if not candidates:
            continue
        best = min(candidates, key=lambda item: _shape_distance(source_shape, item))
        candidates.remove(best)
        pairs.append((source_shape, best))
    return pairs


def _run_sizes(text_frame):
    values = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            values.append(None if run.font.size is None else round(run.font.size.pt, 3))
    return values


def _restore_text_frame_scale(source_tf, output_tf):
    # Source-faithful means the source typography scale is immutable.  The
    # beautifier may improve font family, color and weight, but must not invent
    # larger type or change fitting behaviour and thereby create overflow.
    output_tf.word_wrap = source_tf.word_wrap
    output_tf.auto_size = source_tf.auto_size
    output_tf.vertical_anchor = source_tf.vertical_anchor

    source_runs = [run for paragraph in source_tf.paragraphs for run in paragraph.runs]
    output_runs = [run for paragraph in output_tf.paragraphs for run in paragraph.runs]

    if len(source_runs) == len(output_runs):
        for source_run, output_run in zip(source_runs, output_runs):
            output_run.font.size = source_run.font.size
        return

    # Fail-safe for unusual run segmentation.  Never let the output exceed the
    # largest explicit source size; if the source inherits its size, remove the
    # explicit size introduced by the beautifier.
    explicit = [run.font.size.pt for run in source_runs if run.font.size is not None]
    if not explicit:
        for output_run in output_runs:
            output_run.font.size = None
        return

    ceiling = max(explicit)
    for output_run in output_runs:
        if output_run.font.size is not None and output_run.font.size.pt > ceiling:
            output_run.font.size = Pt(ceiling)


def _restore_table_scale(source_shape, output_shape):
    source = source_shape.table
    output = output_shape.table
    if len(source.rows) != len(output.rows) or len(source.columns) != len(output.columns):
        return
    for r in range(len(source.rows)):
        for c in range(len(source.columns)):
            _restore_text_frame_scale(source.cell(r, c).text_frame, output.cell(r, c).text_frame)


def _restore_source_faithful_typography(source_path: Path, output_path: Path):
    source = Presentation(str(source_path))
    output = Presentation(str(output_path))
    if len(source.slides) != len(output.slides):
        raise RuntimeError("TYPOGRAPHY_GUARD_SLIDE_COUNT_MISMATCH")

    for source_slide, output_slide in zip(source.slides, output.slides):
        for source_shape, output_shape in _match_text_shapes(source_slide, output_slide):
            _restore_text_frame_scale(source_shape.text_frame, output_shape.text_frame)
        for source_shape, output_shape in _match_tables(source_slide, output_slide):
            _restore_table_scale(source_shape, output_shape)

    output.save(str(output_path))


def _new_bounds_violation(source_shape, output_shape, slide_w: int, slide_h: int) -> bool:
    def outside(shape):
        return (
            int(shape.left) < 0
            or int(shape.top) < 0
            or int(shape.left + shape.width) > slide_w
            or int(shape.top + shape.height) > slide_h
        )

    return outside(output_shape) and not outside(source_shape)


def _new_footer_collision(source_shape, output_shape, slide_h: int) -> bool:
    # Bottom 8% is a protected footer band.  Existing source content there is
    # tolerated, but beautification may not newly push content into that band.
    footer_top = int(slide_h * 0.92)
    source_bottom = int(source_shape.top + source_shape.height)
    output_bottom = int(output_shape.top + output_shape.height)
    return source_bottom <= footer_top < output_bottom


def _verify_typography_guard(source_path: Path, output_path: Path):
    source = Presentation(str(source_path))
    output = Presentation(str(output_path))
    violations = []

    for slide_index, (source_slide, output_slide) in enumerate(zip(source.slides, output.slides), 1):
        for source_shape, output_shape in _match_text_shapes(source_slide, output_slide):
            before = _run_sizes(source_shape.text_frame)
            after = _run_sizes(output_shape.text_frame)
            if before != after:
                violations.append(
                    f"slide {slide_index}: source-faithful font-size drift for text={_norm_text(source_shape.text_frame.text)[:80]!r}"
                )
            if _new_bounds_violation(source_shape, output_shape, source.slide_width, source.slide_height):
                violations.append(
                    f"slide {slide_index}: new text-box slide-bounds overflow for text={_norm_text(source_shape.text_frame.text)[:80]!r}"
                )
            if _new_footer_collision(source_shape, output_shape, source.slide_height):
                violations.append(
                    f"slide {slide_index}: new footer collision for text={_norm_text(source_shape.text_frame.text)[:80]!r}"
                )

        for source_shape, output_shape in _match_tables(source_slide, output_slide):
            source_table = source_shape.table
            output_table = output_shape.table
            if len(source_table.rows) != len(output_table.rows) or len(source_table.columns) != len(output_table.columns):
                continue
            for r in range(len(source_table.rows)):
                for c in range(len(source_table.columns)):
                    if _run_sizes(source_table.cell(r, c).text_frame) != _run_sizes(output_table.cell(r, c).text_frame):
                        violations.append(f"slide {slide_index}: table font-size drift at cell {r},{c}")

    return violations


def beautify_pptx(source: str | Path, output: str | Path, style: str, log=lambda _text: None) -> RunReport:
    source_path = Path(source)
    output_path = Path(output)

    # Run the proven v0.7.1 content-preserving engine first, but suppress its
    # success claim until the v0.7.2 typography guard has completed.
    def base_log(message: str):
        if message.strip() == "OFFLINE_BEAUTIFY_PASS=true":
            return
        log(message)

    report = base_engine.beautify_pptx(source_path, output_path, style, base_log)

    if style == SOURCE_FAITHFUL_STYLE:
        _restore_source_faithful_typography(source_path, output_path)
        violations = _verify_typography_guard(source_path, output_path)
        if violations:
            output_path.unlink(missing_ok=True)
            preview = "\n".join(violations[:10])
            raise RuntimeError(f"TYPOGRAPHY_SCALE_GUARD_FAIL: {len(violations)} violations\n{preview}")

        # Formatting repair must remain semantically identical.
        if base_engine.build_manifest is None or base_engine.content_diff is None:
            output_path.unlink(missing_ok=True)
            raise RuntimeError("CONTENT_LOCK_HELPER_MISSING_AFTER_TYPOGRAPHY_REPAIR")
        before = base_engine.build_manifest(str(source_path))
        after = base_engine.build_manifest(str(output_path))
        diffs = base_engine.content_diff(before, after)
        if diffs:
            output_path.unlink(missing_ok=True)
            preview = "\n".join(diffs[:8])
            raise RuntimeError(f"CONTENT_LOCK_FAIL_AFTER_TYPOGRAPHY_REPAIR: {len(diffs)} differences\n{preview}")

        log("SOURCE_FONT_SIZE_LOCK_PASS=true")
        log("TYPOGRAPHY_RATIO_PASS=true")
        log("TEXT_BOUNDS_GUARD_PASS=true")
        log("FOOTER_COLLISION_GUARD_PASS=true")
        log("TYPOGRAPHY_SCALE_GUARD_PASS=true")
        log("CONTENT_LOCK_AFTER_TYPOGRAPHY_PASS=true")
    else:
        log("SOURCE_FONT_SIZE_LOCK_PASS=not_applicable")
        log("TYPOGRAPHY_SCALE_GUARD_PASS=not_applicable")

    log(f"OFFLINE_ENGINE_VERSION={ENGINE_VERSION}")
    log("OFFLINE_BEAUTIFY_PASS=true")
    return report
