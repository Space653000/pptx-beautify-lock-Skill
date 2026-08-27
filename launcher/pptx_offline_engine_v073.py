from __future__ import annotations

from pathlib import Path
import os
import shutil
import tempfile
import zipfile

from lxml import etree
from pptx import Presentation

import pptx_offline_engine as base_engine

ENGINE_VERSION = "0.7.3"
STYLE_PRESETS = base_engine.STYLE_PRESETS
RunReport = base_engine.RunReport
SOURCE_FAITHFUL_STYLE = "自動（忠於原稿 / Source-faithful）"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PROOFING_NODES = ("rPr", "defRPr", "endParaRPr")


def _is_slide_xml(name: str) -> bool:
    return name.startswith("ppt/slides/slide") and name.endswith(".xml")


def _set_no_proof(xml_bytes: bytes) -> tuple[bytes, int]:
    """Set editor proofing metadata without changing visual layout/style.

    Source-faithful v0.7.3 intentionally does not create new run-property nodes.
    It only adds ``noProof=1`` to existing DrawingML run/default/end-paragraph
    property nodes.  This keeps the change budget limited to non-rendering
    proofing metadata.
    """
    parser = etree.XMLParser(remove_blank_text=False, resolve_entities=False)
    root = etree.fromstring(xml_bytes, parser=parser)
    changed = 0
    for local_name in PROOFING_NODES:
        for node in root.findall(f".//{{{A_NS}}}{local_name}"):
            if node.get("noProof") != "1":
                node.set("noProof", "1")
                changed += 1
    if not changed:
        return xml_bytes, 0
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=None), changed


def _rewrite_proofing_metadata(source_path: Path, output_path: Path) -> int:
    """Copy the PPTX and modify slide proofing metadata only."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_path, output_path)

    fd, temp_name = tempfile.mkstemp(
        prefix=f".{output_path.stem}.proofing-",
        suffix=".pptx",
        dir=output_path.parent,
    )
    os.close(fd)
    temp_path = Path(temp_name)
    changed = 0
    try:
        with zipfile.ZipFile(output_path, "r") as zin, zipfile.ZipFile(temp_path, "w") as zout:
            for info in zin.infolist():
                payload = zin.read(info.filename)
                if _is_slide_xml(info.filename):
                    payload, delta = _set_no_proof(payload)
                    changed += delta
                zout.writestr(info, payload)
        os.replace(temp_path, output_path)
    finally:
        temp_path.unlink(missing_ok=True)
    return changed


def _canonical_without_no_proof(xml_bytes: bytes) -> bytes:
    parser = etree.XMLParser(remove_blank_text=False, resolve_entities=False)
    root = etree.fromstring(xml_bytes, parser=parser)
    for node in root.xpath("//*[@noProof]"):
        node.attrib.pop("noProof", None)
    return etree.tostring(root, method="c14n", with_comments=True)


def _verify_only_allowlisted_package_change(source_path: Path, output_path: Path) -> list[str]:
    """Prove that Source-faithful changed only slide ``noProof`` metadata."""
    violations: list[str] = []
    with zipfile.ZipFile(source_path, "r") as source_zip, zipfile.ZipFile(output_path, "r") as output_zip:
        source_names = source_zip.namelist()
        output_names = output_zip.namelist()
        if source_names != output_names:
            violations.append("package member/order drift")
            return violations

        for name in source_names:
            before = source_zip.read(name)
            after = output_zip.read(name)
            if before == after:
                continue
            if not _is_slide_xml(name):
                violations.append(f"non-slide package part changed: {name}")
                continue
            try:
                if _canonical_without_no_proof(before) != _canonical_without_no_proof(after):
                    violations.append(f"slide XML changed beyond proofing metadata: {name}")
            except Exception as exc:
                violations.append(f"cannot verify slide XML {name}: {type(exc).__name__}")
    return violations


def _verify_content_lock(source_path: Path, output_path: Path) -> list[str]:
    if base_engine.build_manifest is None or base_engine.content_diff is None:
        return ["CONTENT_LOCK_HELPER_MISSING"]
    before = base_engine.build_manifest(str(source_path))
    after = base_engine.build_manifest(str(output_path))
    return list(base_engine.content_diff(before, after))


def _source_faithful_safe_only(
    source_path: Path,
    output_path: Path,
    log,
) -> RunReport:
    if source_path.resolve() == output_path.resolve():
        raise ValueError("output must not overwrite source")

    slide_count = len(Presentation(str(source_path)).slides)
    proofing_nodes = _rewrite_proofing_metadata(source_path, output_path)

    package_violations = _verify_only_allowlisted_package_change(source_path, output_path)
    if package_violations:
        output_path.unlink(missing_ok=True)
        preview = "\n".join(package_violations[:10])
        raise RuntimeError(
            f"SOURCE_FAITHFUL_PACKAGE_GUARD_FAIL: {len(package_violations)} violations\n{preview}"
        )

    content_diffs = _verify_content_lock(source_path, output_path)
    if content_diffs:
        output_path.unlink(missing_ok=True)
        preview = "\n".join(content_diffs[:8])
        raise RuntimeError(
            f"CONTENT_LOCK_FAIL_AFTER_SAFE_ONLY: {len(content_diffs)} differences\n{preview}"
        )

    # Because every PPTX package part is byte-identical except noProof attributes
    # inside slide XML, these visual contracts are stronger than a heuristic
    # geometry comparison: source geometry, typography, table style, media and
    # theme/master identity are unchanged by construction.
    log("SOURCE_FAITHFUL_SAFE_ONLY=true")
    log("SOURCE_CHANGE_POLICY=proofing_metadata_only")
    log(f"PROOFING_METADATA_NODES_UPDATED={proofing_nodes}")
    log("SOURCE_PACKAGE_STRUCTURE_PASS=true")
    log("SOURCE_VISUAL_XML_LOCK_PASS=true")
    log("SOURCE_GEOMETRY_LOCK_PASS=true")
    log("SOURCE_TYPOGRAPHY_LOCK_PASS=true")
    log("SOURCE_TABLE_STYLE_LOCK_PASS=true")
    log("SOURCE_MEDIA_LOCK_PASS=true")
    log("SOURCE_THEME_IDENTITY_LOCK_PASS=true")
    log("SAFE_CHANGE_BUDGET_PASS=true")
    log("NO_DEGRADATION_GATE_PASS=true")
    log("CONTENT_LOCK_PASS=true")

    return RunReport(
        slide_count=slide_count,
        removed_empty_placeholders=0,
        suppressed_template_artifacts=0,
        tables_styled=0,
        data_slides_structured=0,
        text_runs_normalized=0,
        warnings=[
            "Source-faithful safe-only mode preserved source visual geometry/style; "
            "only proofing metadata was allowlisted."
        ],
    )


def beautify_pptx(source: str | Path, output: str | Path, style: str, log=lambda _text: None) -> RunReport:
    source_path = Path(source)
    output_path = Path(output)

    if style == SOURCE_FAITHFUL_STYLE:
        report = _source_faithful_safe_only(source_path, output_path, log)
    else:
        # Transformative presets remain explicit opt-in.  They use the existing
        # deterministic formatter, but are intentionally NOT labelled
        # Source-faithful and do not receive the safe-only/no-degradation claim.
        def delegated_log(message: str):
            stripped = message.strip()
            if stripped == "OFFLINE_BEAUTIFY_PASS=true" or stripped.startswith("OFFLINE_ENGINE_VERSION="):
                return
            log(message)

        report = base_engine.beautify_pptx(source_path, output_path, style, delegated_log)
        log("SOURCE_FAITHFUL_SAFE_ONLY=not_applicable")
        log("SAFE_CHANGE_BUDGET_PASS=not_applicable_transformative_style")
        log("NO_DEGRADATION_GATE_PASS=not_claimed_transformative_style")

    log(f"OFFLINE_ENGINE_VERSION={ENGINE_VERSION}")
    log("OFFLINE_BEAUTIFY_PASS=true")
    return report
