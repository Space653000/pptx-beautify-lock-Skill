from __future__ import annotations

from pathlib import Path
import sys
import tempfile
import unittest
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
for path in (ROOT / "launcher", ROOT / "pptx-beautify-lock" / "scripts"):
    value = str(path)
    if value not in sys.path:
        sys.path.insert(0, value)

import pptx_offline_engine_v073 as engine

STYLE = "自動（忠於原稿 / Source-faithful）"


def _shape_fingerprint(prs: Presentation):
    slides = []
    for slide in prs.slides:
        shapes = []
        for shape in slide.shapes:
            item = {
                "type": int(shape.shape_type),
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
                "rotation": float(shape.rotation or 0),
                "placeholder": bool(getattr(shape, "is_placeholder", False)),
                "text": getattr(shape, "text", ""),
            }
            if getattr(shape, "has_text_frame", False):
                item["word_wrap"] = shape.text_frame.word_wrap
                item["auto_size"] = shape.text_frame.auto_size
                item["vertical_anchor"] = shape.text_frame.vertical_anchor
                item["runs"] = [
                    (
                        run.text,
                        run.font.name,
                        None if run.font.size is None else round(run.font.size.pt, 3),
                        run.font.bold,
                        run.font.italic,
                        run.font.underline,
                    )
                    for p in shape.text_frame.paragraphs
                    for run in p.runs
                ]
            if getattr(shape, "has_table", False):
                table = shape.table
                item["table"] = {
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                    "row_heights": [int(row.height) for row in table.rows],
                    "col_widths": [int(col.width) for col in table.columns],
                    "text": [
                        [table.cell(r, c).text for c in range(len(table.columns))]
                        for r in range(len(table.rows))
                    ],
                }
            if shape.shape_type == 13:  # picture
                item["crop"] = (
                    shape.crop_left,
                    shape.crop_right,
                    shape.crop_top,
                    shape.crop_bottom,
                )
                item["image_sha1"] = shape.image.sha1
            shapes.append(item)
        slides.append(shapes)
    return slides


class SourceFaithfulSafeOnlyV073Tests(unittest.TestCase):
    def test_source_faithful_preserves_visual_fingerprint(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            src = root / "source.pptx"
            out = root / "out.pptx"
            image_path = root / "image.png"
            Image.new("RGB", (320, 180), (230, 235, 240)).save(image_path)

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "PMX Acoustic Lab Construction Assessment Report"
            # Leave subtitle placeholder empty intentionally. Safe-only mode must
            # not delete/reposition it just to make the file look changed.
            title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
            title_run.font.name = "Arial"
            title_run.font.size = Pt(23)
            slide.shapes.add_picture(str(image_path), Inches(7.0), Inches(2.0), width=Inches(2.5))

            refs = prs.slides.add_slide(prs.slide_layouts[6])
            box = refs.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(5.5))
            tf = box.text_frame
            lines = [
                ("Sources & References (1/4)", 18),
                ("Item 4 - APX500 software V10.0", 20),
                ("URL: https://www.ap.com/analyzers-accessories/apx52x", 16),
                ("AP official product page and technical library", 13),
            ]
            tf.clear()
            for index, (text, size) in enumerate(lines):
                p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
                r = p.add_run()
                r.text = text
                r.font.name = "Arial"
                r.font.size = Pt(size)

            table_shape = refs.shapes.add_table(3, 3, Inches(0.7), Inches(5.8), Inches(7.0), Inches(1.0))
            for r in range(3):
                for c in range(3):
                    cell = table_shape.table.cell(r, c)
                    cell.text = f"R{r}C{c}"
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(9 + r + c)
            prs.save(src)

            before = _shape_fingerprint(Presentation(src))
            logs: list[str] = []
            report = engine.beautify_pptx(src, out, STYLE, logs.append)
            after = _shape_fingerprint(Presentation(out))

            self.assertEqual(before, after)
            self.assertEqual(report.removed_empty_placeholders, 0)
            self.assertEqual(report.tables_styled, 0)
            self.assertEqual(report.data_slides_structured, 0)
            self.assertIn("SOURCE_FAITHFUL_SAFE_ONLY=true", logs)
            self.assertIn("SOURCE_CHANGE_POLICY=proofing_metadata_only", logs)
            self.assertIn("SOURCE_PACKAGE_STRUCTURE_PASS=true", logs)
            self.assertIn("SOURCE_VISUAL_XML_LOCK_PASS=true", logs)
            self.assertIn("SOURCE_GEOMETRY_LOCK_PASS=true", logs)
            self.assertIn("SOURCE_TYPOGRAPHY_LOCK_PASS=true", logs)
            self.assertIn("SOURCE_TABLE_STYLE_LOCK_PASS=true", logs)
            self.assertIn("SOURCE_MEDIA_LOCK_PASS=true", logs)
            self.assertIn("SOURCE_THEME_IDENTITY_LOCK_PASS=true", logs)
            self.assertIn("SAFE_CHANGE_BUDGET_PASS=true", logs)
            self.assertIn("NO_DEGRADATION_GATE_PASS=true", logs)
            self.assertIn("CONTENT_LOCK_PASS=true", logs)
            self.assertIn("OFFLINE_ENGINE_VERSION=0.7.3", logs)

    def test_package_change_is_limited_to_slide_proofing_metadata(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            src = root / "source.pptx"
            out = root / "out.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Reference-heavy slide should not be redesigned"
            prs.save(src)

            engine.beautify_pptx(src, out, STYLE)
            self.assertEqual(engine._verify_only_allowlisted_package_change(src, out), [])

            with zipfile.ZipFile(src, "r") as a, zipfile.ZipFile(out, "r") as b:
                self.assertEqual(a.namelist(), b.namelist())
                changed = [name for name in a.namelist() if a.read(name) != b.read(name)]
                self.assertTrue(all(name.startswith("ppt/slides/slide") for name in changed))

    def test_wrapper_declares_v073_contract(self):
        self.assertEqual(engine.ENGINE_VERSION, "0.7.3")
        self.assertTrue(callable(engine.beautify_pptx))
        self.assertTrue(callable(engine._verify_only_allowlisted_package_change))


if __name__ == "__main__":
    unittest.main()
