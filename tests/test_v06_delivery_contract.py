from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REGRESSION = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_regression.py"


def run_regression(*args: str):
    return subprocess.run(
        [sys.executable, str(REGRESSION), *map(str, args)],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


def make_clean_deck(path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(6.0), Inches(0.8))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "v0.6 Global Design Jury / 世界級交付契約"
    run.font.size = Pt(24)
    prs.save(path)


def visual_report(path: Path):
    checks = {
        "no_unintended_overlap": True,
        "no_clipping_or_overflow": True,
        "content_visible": True,
        "text_readable": True,
        "hierarchy_clear": True,
        "alignment_consistent": True,
        "tables_charts_readable": True,
        "style_consistent": True,
        "no_template_placeholder_artifacts": True,
        "theme_fidelity_preserved": True,
        "bilingual_typography_clean": True,
    }
    path.write_text(
        json.dumps(
            {
                "schema": 3,
                "slide_count": 1,
                "render_engine": "contract-test-renderer",
                "reviewer": "contract-test-reviewer",
                "overall_pass": True,
                "slides": [{"slide": 1, "score": 96, "checks": checks}],
            }
        ),
        encoding="utf-8",
    )


def composition_report(path: Path):
    checks = {
        "brand_chrome_respected": True,
        "content_not_occluded": True,
        "grid_alignment_coherent": True,
        "peer_components_aligned": True,
        "spacing_rhythm_coherent": True,
        "reading_order_clear": True,
        "visual_balance_coherent": True,
        "slide_role_composition_fit": True,
        "decorative_elements_earn_their_place": True,
    }
    scores = {
        "hierarchy": 94,
        "alignment": 94,
        "spacing": 94,
        "balance": 93,
        "brand_fidelity": 96,
        "restraint": 94,
        "data_legibility": 94,
    }
    path.write_text(
        json.dumps(
            {
                "schema": 1,
                "slide_count": 1,
                "render_engine": "contract-test-renderer",
                "reviewer": "contract-test-reviewer",
                "overall_pass": True,
                "slides": [
                    {
                        "slide": 1,
                        "composition_score": 94,
                        "checks": checks,
                        "scores": scores,
                        "evidence": {
                            "source_comparison": "source and final compared",
                            "grid_rails": ["title-left", "content-top"],
                            "reading_order": ["title"],
                            "brand_anchors": ["synthetic source identity retained"],
                        },
                    }
                ],
            }
        ),
        encoding="utf-8",
    )


def jury_report(path: Path, identity_score: float = 97):
    checks = {
        "purpose_is_clear": True,
        "focal_point_is_unambiguous": True,
        "hierarchy_is_structural": True,
        "spacing_is_intentional": True,
        "typography_is_crafted": True,
        "color_is_disciplined": True,
        "source_identity_is_preserved": True,
        "signal_to_noise_is_high": True,
        "glance_test_pass": True,
        "brand_and_status_do_not_compete": True,
        "no_generic_template_skin": True,
    }
    scores = {
        "purpose": 94,
        "hierarchy": 94,
        "simplicity": 94,
        "craft": 94,
        "composition": 94,
        "typography": 94,
        "spacing_rhythm": 94,
        "color_discipline": 94,
        "source_identity": identity_score,
        "signal_to_noise": 94,
        "glance_readability": 94,
        "executive_readiness": 94,
    }
    role_scores = {
        "role_fit": 94,
        "audience_fit": 94,
        "information_density_control": 94,
        "reading_path": 94,
        "visual_coherence": 94,
    }
    payload = {
        "schema": 1,
        "slide_count": 1,
        "render_engine": "contract-test-renderer",
        "reviewer": "contract-test-reviewer",
        "audience_profile": "world-class external technology customer",
        "review_rounds": 2,
        "overall_pass": True,
        "deck_jury_score": 94,
        "jury_lenses": {
            "purpose_hierarchy_craft": {"pass": True, "evidence": "craft lens pass"},
            "executive_communication": {"pass": True, "evidence": "communication lens pass"},
            "domain_role_fit": {"pass": True, "evidence": "role-fit lens pass"},
        },
        "deck_identity": {
            "checks": {
                "source_personality_preserved": True,
                "no_template_convergence": True,
                "no_unjustified_cardification": True,
                "no_unjustified_dark_techification": True,
                "no_unjustified_gradientization": True,
                "no_brand_personality_erasure": True,
            },
            "source_personality": "synthetic source",
            "final_personality": "same source with refined craft",
            "identity_evidence": ["same canvas", "same content identity"],
            "identity_fidelity_score": identity_score,
            "archetype_fit_score": 94,
            "generic_template_risk": 3,
        },
        "slides": [
            {
                "slide": 1,
                "jury_role": "other",
                "checks": checks,
                "scores": scores,
                "slide_jury_score": 94,
                "role_scores": role_scores,
                "role_score": 94,
                "evidence": {
                    "primary_purpose": "validate strict delivery contract",
                    "focal_point": "single title",
                    "reading_order": ["title"],
                    "grid_or_alignment_logic": ["single left rail"],
                    "spacing_logic": "intentional single-object spacing",
                    "source_identity_anchors": ["source title and canvas"],
                    "what_was_removed_or_restrained": "nothing unnecessary added",
                    "why_this_is_not_a_generic_template": "source geometry remains authoritative",
                },
            }
        ],
    }
    path.write_text(json.dumps(payload), encoding="utf-8")


class V06DeliveryContractTests(unittest.TestCase):
    def test_v06_delivery_requires_and_passes_all_three_render_jury_reports(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "output.pptx"
            visual = Path(td) / "visual.json"
            composition = Path(td) / "composition.json"
            jury = Path(td) / "jury.json"
            make_clean_deck(source)
            Presentation(source).save(output)
            visual_report(visual)
            composition_report(composition)
            jury_report(jury)

            result = run_regression(
                source,
                output,
                "--visual-qa-report",
                visual,
                "--require-visual-qa",
                "--composition-qa-report",
                composition,
                "--require-composition-qa",
                "--global-jury-report",
                jury,
                "--require-global-jury",
            )
            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            self.assertIn("CONTENT_LOCK_PASS=true", result.stdout)
            self.assertIn("COMPOSITION_QA_PASS=true", result.stdout)
            self.assertIn("DECK_IDENTITY_PASS=true", result.stdout)
            self.assertIn("GLOBAL_DESIGN_JURY_PASS=true", result.stdout)
            self.assertIn("REGRESSION_V06_PASS=true", result.stdout)
            self.assertIn("DELIVERY_V06_PASS=true", result.stdout)

    def test_v06_fails_closed_without_global_jury_report(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "output.pptx"
            visual = Path(td) / "visual.json"
            composition = Path(td) / "composition.json"
            make_clean_deck(source)
            Presentation(source).save(output)
            visual_report(visual)
            composition_report(composition)

            result = run_regression(
                source,
                output,
                "--visual-qa-report",
                visual,
                "--require-visual-qa",
                "--composition-qa-report",
                composition,
                "--require-composition-qa",
                "--require-global-jury",
            )
            self.assertNotEqual(0, result.returncode)
            self.assertIn("GLOBAL_DESIGN_JURY_PASS=false", result.stdout)
            self.assertIn("DELIVERY_V06_PASS=false", result.stdout)

    def test_v06_fails_when_source_identity_is_below_world_class_floor(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "output.pptx"
            visual = Path(td) / "visual.json"
            composition = Path(td) / "composition.json"
            jury = Path(td) / "jury.json"
            make_clean_deck(source)
            Presentation(source).save(output)
            visual_report(visual)
            composition_report(composition)
            jury_report(jury, identity_score=94)

            result = run_regression(
                source,
                output,
                "--visual-qa-report",
                visual,
                "--require-visual-qa",
                "--composition-qa-report",
                composition,
                "--require-composition-qa",
                "--global-jury-report",
                jury,
                "--require-global-jury",
            )
            self.assertNotEqual(0, result.returncode)
            self.assertIn("DECK_IDENTITY_PASS=false", result.stdout)
            self.assertIn("GLOBAL_DESIGN_JURY_PASS=false", result.stdout)
            self.assertIn("DELIVERY_V06_PASS=false", result.stdout)


if __name__ == "__main__":
    unittest.main()
