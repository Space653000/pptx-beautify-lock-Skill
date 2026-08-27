from __future__ import annotations

from pathlib import Path
import sys
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
for path in (ROOT / "launcher", ROOT / "pptx-beautify-lock" / "scripts"):
    value = str(path)
    if value not in sys.path:
        sys.path.insert(0, value)

import pptx_offline_engine_v072 as engine

STYLE = "自動（忠於原稿 / Source-faithful）"


def _add_text(slide, text, left, top, width, height, size):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    return shape


def _sizes(prs):
    values = {}
    for slide_index, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            values[(slide_index, text)] = [
                None if run.font.size is None else round(run.font.size.pt, 2)
                for p in shape.text_frame.paragraphs
                for run in p.runs
            ]
    return values


class TypographyScaleGuardV072Tests(unittest.TestCase):
    def test_source_faithful_never_inflates_existing_text(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            src = root / "source.pptx"
            out = root / "out.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_text(slide, "Sources & References (1/4)", 0.5, 0.25, 7.5, 0.6, 18)
            _add_text(slide, "Item 4 - APX500 software V10.0", 0.5, 2.0, 8.0, 0.5, 24)
            _add_text(slide, "URL: https://www.ap.com/analyzers-accessories/apx52x", 0.8, 2.65, 8.5, 0.55, 24)
            _add_text(slide, "Normal explanatory body text that must remain readable.", 0.5, 4.0, 8.0, 0.5, 14)
            prs.save(src)

            before = _sizes(Presentation(src))
            logs = []
            engine.beautify_pptx(src, out, STYLE, logs.append)
            after = _sizes(Presentation(out))

            self.assertEqual(before, after)
            self.assertIn("SOURCE_FONT_SIZE_LOCK_PASS=true", logs)
            self.assertIn("TYPOGRAPHY_RATIO_PASS=true", logs)
            self.assertIn("TEXT_BOUNDS_GUARD_PASS=true", logs)
            self.assertIn("FOOTER_COLLISION_GUARD_PASS=true", logs)
            self.assertIn("TYPOGRAPHY_SCALE_GUARD_PASS=true", logs)
            self.assertIn("CONTENT_LOCK_AFTER_TYPOGRAPHY_PASS=true", logs)

    def test_table_font_sizes_are_preserved_in_source_faithful_mode(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            src = root / "source.pptx"
            out = root / "out.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(7), Inches(2))
            table = table_shape.table
            for r in range(2):
                for c in range(2):
                    cell = table.cell(r, c)
                    cell.text = f"R{r}C{c}"
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(11 + r + c)
            prs.save(src)

            source = Presentation(src)
            source_sizes = [
                source.slides[0].shapes[0].table.cell(r, c).text_frame.paragraphs[0].runs[0].font.size.pt
                for r in range(2) for c in range(2)
            ]
            engine.beautify_pptx(src, out, STYLE)
            result = Presentation(out)
            output_sizes = [
                result.slides[0].shapes[0].table.cell(r, c).text_frame.paragraphs[0].runs[0].font.size.pt
                for r in range(2) for c in range(2)
            ]
            self.assertEqual(source_sizes, output_sizes)

    def test_wrapper_declares_v072_and_contract_surface(self):
        self.assertEqual(engine.ENGINE_VERSION, "0.7.2")
        self.assertTrue(hasattr(engine, "STYLE_PRESETS"))
        self.assertTrue(callable(engine.beautify_pptx))


if __name__ == "__main__":
    unittest.main()
