from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LOCK = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_content_lock.py"
LINT = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_lint.py"
THEME = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_theme_profile.py"
VISUAL_QA = ROOT / "pptx-beautify-lock" / "scripts" / "visual_qa_gate.py"
REGRESSION = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_regression.py"


def run_script(script: Path, *args: str):
    return subprocess.run(
        [sys.executable, str(script), *map(str, args)],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


def visual_checks():
    return {
        "no_unintended_overlap": True,
        "no_clipping_or_overflow": True,
        "content_visible": True,
        "text_readable": True,
        "hierarchy_clear": True,
        "alignment_consistent": True,
        "tables_charts_readable": True,
        "style_consistent": True,
        "no_template_placeholder_artifacts": True,
        "theme_fidelity_preserved": True,
        "bilingual_typography_clean": True,
    }


class ContentLockContractTests(unittest.TestCase):
    def make_source(self, path: Path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "內容不可修改 / Content must not change 123.45%"
        run.font.size = Pt(18)

        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.2), Inches(5), Inches(1.5))
        table = table_shape.table
        table.cell(0, 0).text = "項目"
        table.cell(0, 1).text = "數值"
        table.cell(1, 0).text = "A"
        table.cell(1, 1).text = "42"
        prs.save(path)

    def assert_content_passes(self, src: Path, out: Path):
        result = run_script(LOCK, "verify", src, out)
        self.assertEqual(0, result.returncode, result.stdout + result.stderr)
        self.assertIn("CONTENT_LOCK_PASS=true", result.stdout)

    def assert_content_fails(self, src: Path, out: Path):
        result = run_script(LOCK, "verify", src, out)
        self.assertNotEqual(0, result.returncode, result.stdout + result.stderr)
        self.assertIn("CONTENT_LOCK_PASS=false", result.stdout)

    def test_visual_only_change_passes_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "src.pptx"
            out = Path(td) / "visual_only.pptx"
            self.make_source(src)

            prs = Presentation(src)
            text_shape = prs.slides[0].shapes[0]
            for p in text_shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(30)
                    r.font.bold = True
            text_shape.left = Inches(1.5)
            prs.save(out)
            self.assert_content_passes(src, out)

    def test_run_segmentation_change_passes_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "one-run.pptx"
            out = Path(td) / "two-runs.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            box.text_frame.paragraphs[0].add_run().text = "Hello world"
            prs.save(src)

            prs2 = Presentation(src)
            p = prs2.slides[0].shapes[0].text_frame.paragraphs[0]
            p.clear()
            r1 = p.add_run(); r1.text = "Hello "; r1.font.bold = True
            r2 = p.add_run(); r2.text = "world"; r2.font.italic = True
            prs2.save(out)
            self.assert_content_passes(src, out)

    def test_text_change_fails_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "src.pptx"; out = Path(td) / "changed.pptx"
            self.make_source(src)
            prs = Presentation(src)
            prs.slides[0].shapes[0].text = "內容被修改 / Content changed"
            prs.save(out)
            self.assert_content_fails(src, out)

    def test_table_value_change_fails_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "src.pptx"; out = Path(td) / "changed_table.pptx"
            self.make_source(src)
            prs = Presentation(src)
            prs.slides[0].shapes[1].table.cell(1, 1).text = "43"
            prs.save(out)
            self.assert_content_fails(src, out)

    def test_table_merge_semantics_change_fails_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "merged.pptx"; out = Path(td) / "split.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2)).table
            table.cell(0, 0).text = "Merged"
            table.cell(0, 1).text = ""
            table.cell(1, 0).text = "A"
            table.cell(1, 1).text = "B"
            table.cell(0, 0).merge(table.cell(0, 1))
            prs.save(src)
            prs2 = Presentation(src)
            prs2.slides[0].shapes[0].table.cell(0, 0).split()
            prs2.save(out)
            self.assert_content_fails(src, out)

    def test_hyperlink_target_change_fails_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "link-a.pptx"; out = Path(td) / "link-b.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            run = box.text_frame.paragraphs[0].add_run()
            run.text = "Open link"; run.hyperlink.address = "https://example.com/a"
            prs.save(src)
            prs2 = Presentation(src)
            prs2.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].hyperlink.address = "https://example.com/b"
            prs2.save(out)
            self.assert_content_fails(src, out)

    def test_swapping_hyperlinks_between_objects_fails_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "links-source.pptx"; out = Path(td) / "links-swapped.pptx"
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
            for index, (label, target) in enumerate((("Document A", "https://example.com/a"), ("Document B", "https://example.com/b"))):
                box = slide.shapes.add_textbox(Inches(1), Inches(1 + index), Inches(5), Inches(0.7))
                run = box.text_frame.paragraphs[0].add_run(); run.text = label; run.hyperlink.address = target
            prs.save(src)
            prs2 = Presentation(src)
            a_run = prs2.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
            b_run = prs2.slides[0].shapes[1].text_frame.paragraphs[0].runs[0]
            a_run.hyperlink.address = "https://example.com/b"
            b_run.hyperlink.address = "https://example.com/a"
            prs2.save(out)
            self.assert_content_fails(src, out)

    def test_hidden_slide_state_change_fails_content_lock(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "hidden.pptx"; out = Path(td) / "shown.pptx"
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Hidden slide"
            slide._element.set("show", "0"); prs.save(src)
            prs2 = Presentation(src); prs2.slides[0]._element.set("show", "1"); prs2.save(out)
            self.assert_content_fails(src, out)

    def test_linter_public_cli_parses_valid_deck(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "src.pptx"; self.make_source(src)
            result = run_script(LINT, src, "--json")
            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            payload = json.loads(result.stdout)
            self.assertEqual(1, payload["slides_checked"])

    def test_linter_detects_tiny_text_inside_table(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "tiny-table.pptx"; self.make_source(src)
            prs = Presentation(src)
            cell = prs.slides[0].shapes[1].table.cell(1, 1)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(6)
            prs.save(src)
            result = run_script(LINT, src, "--json")
            payload = json.loads(result.stdout)
            self.assertIn("tiny-text", [f["rule"] for f in payload["findings"]])

    def test_linter_detects_placeholder_artifact(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "placeholder.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "presentation title"
            prs.save(src)
            result = run_script(LINT, src, "--json")
            self.assertNotEqual(0, result.returncode)
            payload = json.loads(result.stdout)
            self.assertIn("template-placeholder-artifact", [f["rule"] for f in payload["findings"]])

    def test_linter_detects_cjk_font_fallback_risk(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "font-risk.pptx"
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            run = box.text_frame.paragraphs[0].add_run(); run.text = "繁體中文 English"; run.font.name = "Inter"
            prs.save(src)
            result = run_script(LINT, src, "--json")
            payload = json.loads(result.stdout)
            self.assertIn("cjk-font-fallback-risk", [f["rule"] for f in payload["findings"]])

    def test_theme_guard_blocks_light_to_dark_inversion(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "light.pptx"; out = Path(td) / "dark.pptx"
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Source"
            prs.save(src)
            prs2 = Presentation(src); s = prs2.slides[0]
            bg = s.shapes.add_shape(1, 0, 0, prs2.slide_width, prs2.slide_height)
            bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(10, 25, 45); bg.line.fill.background()
            # Move background behind all other shapes by XML order.
            sp_tree = s.shapes._spTree
            sp_tree.remove(bg._element); sp_tree.insert(2, bg._element)
            prs2.save(out)
            result = run_script(THEME, "compare", src, out, "--json")
            self.assertNotEqual(0, result.returncode, result.stdout + result.stderr)
            payload = json.loads(result.stdout)
            self.assertFalse(payload["THEME_GUARD_PASS"])

    def test_theme_guard_allows_light_source_with_small_accent(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "light.pptx"; out = Path(td) / "accent.pptx"
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Source"
            prs.save(src)
            prs2 = Presentation(src); s = prs2.slides[0]
            accent = s.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(0.2), Inches(5))
            accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(0, 120, 160); accent.line.fill.background()
            prs2.save(out)
            result = run_script(THEME, "compare", src, out, "--json")
            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            self.assertTrue(json.loads(result.stdout)["THEME_GUARD_PASS"])

    def test_visual_qa_requires_every_slide_and_new_checks(self):
        with tempfile.TemporaryDirectory() as td:
            report_path = Path(td) / "visual_qa.json"
            report = {
                "schema": 3,
                "slide_count": 2,
                "render_engine": "PowerPoint",
                "reviewer": "AI vision reviewer",
                "overall_pass": True,
                "slides": [
                    {"slide": 1, "score": 90, "checks": visual_checks()},
                    {"slide": 2, "score": 88, "checks": visual_checks()},
                ],
            }
            report_path.write_text(json.dumps(report), encoding="utf-8")
            good = run_script(VISUAL_QA, report_path, "--expected-slides", "2")
            self.assertEqual(0, good.returncode, good.stdout + good.stderr)

            report["slides"][1]["checks"]["theme_fidelity_preserved"] = False
            report_path.write_text(json.dumps(report), encoding="utf-8")
            bad = run_script(VISUAL_QA, report_path, "--expected-slides", "2")
            self.assertNotEqual(0, bad.returncode)
            self.assertIn("VISUAL_QA_PASS=false", bad.stdout)

    def test_regression_requires_visual_report_for_delivery(self):
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "src.pptx"; out = Path(td) / "out.pptx"
            self.make_source(src); Presentation(src).save(out)
            result = run_script(REGRESSION, src, out, "--require-visual-qa")
            self.assertNotEqual(0, result.returncode)
            self.assertIn("DELIVERY_PASS=false", result.stdout)
            self.assertIn("VISUAL_QA_PASS=false", result.stdout)

    def test_plugin_manifests_point_to_installable_skill(self):
        plugin = json.loads((ROOT / ".claude-plugin" / "plugin.json").read_text(encoding="utf-8"))
        marketplace = json.loads((ROOT / ".claude-plugin" / "marketplace.json").read_text(encoding="utf-8"))
        self.assertEqual("pptx-beautify-lock", plugin["name"])
        self.assertIn("./pptx-beautify-lock", plugin["skills"])
        self.assertTrue((ROOT / "pptx-beautify-lock" / "SKILL.md").exists())
        self.assertEqual("space653000-pptx", marketplace["name"])
        self.assertEqual(plugin["name"], marketplace["plugins"][0]["name"])


if __name__ == "__main__":
    unittest.main()
