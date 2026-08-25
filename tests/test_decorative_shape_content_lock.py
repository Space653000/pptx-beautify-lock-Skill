from __future__ import annotations

import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
LOCK = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_content_lock.py"


def run_lock(source: Path, output: Path):
    return subprocess.run(
        [sys.executable, str(LOCK), "verify", str(source), str(output)],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


class DecorativeShapeContentLockTests(unittest.TestCase):
    def make_source(self, path: Path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Protected content / 受保護內容"
        prs.save(path)

    def test_content_free_native_decorative_shape_is_visual_only(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "decorated.pptx"
            self.make_source(source)

            prs = Presentation(source)
            slide = prs.slides[0]
            shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(0.2), Inches(5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 120, 160)
            shape.line.fill.background()
            prs.save(output)

            result = run_lock(source, output)
            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            self.assertIn("CONTENT_LOCK_PASS=true", result.stdout)
            self.assertIn("content_differences=0", result.stdout)

    def test_new_visible_text_inside_shape_is_still_content_change(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "new-text.pptx"
            self.make_source(source)

            prs = Presentation(source)
            slide = prs.slides[0]
            shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(3), Inches(1))
            shape.text = "NEW CONTENT"
            prs.save(output)

            result = run_lock(source, output)
            self.assertNotEqual(0, result.returncode)
            self.assertIn("CONTENT_LOCK_PASS=false", result.stdout)


if __name__ == "__main__":
    unittest.main()
