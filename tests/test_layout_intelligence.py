from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
SPATIAL = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_layout_intelligence.py"
COMPOSITION = ROOT / "pptx-beautify-lock" / "scripts" / "composition_qa_gate.py"


def run_script(script: Path, *args: str):
    return subprocess.run(
        [sys.executable, str(script), *map(str, args)],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


def make_source(path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Protected title"
    prs.save(path)


def composition_checks():
    return {
        "brand_chrome_respected": True,
        "content_not_occluded": True,
        "grid_alignment_coherent": True,
        "peer_components_aligned": True,
        "spacing_rhythm_coherent": True,
        "reading_order_clear": True,
        "visual_balance_coherent": True,
        "slide_role_composition_fit": True,
        "decorative_elements_earn_their_place": True,
    }


def composition_scores(value=92):
    return {
        "hierarchy": value,
        "alignment": value,
        "spacing": value,
        "balance": value,
        "brand_fidelity": value,
        "restraint": value,
        "data_legibility": value,
    }


class LayoutIntelligenceTests(unittest.TestCase):
    def test_foreground_solid_fill_over_text_fails_spatial_gate(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "bad.pptx"
            make_source(source)

            prs = Presentation(source)
            slide = prs.slides[0]
            cover = slide.shapes.add_shape(1, Inches(1.2), Inches(1.1), Inches(4), Inches(0.8))
            cover.fill.solid()
            cover.fill.fore_color.rgb = RGBColor(0, 0, 0)
            cover.line.fill.background()
            prs.save(output)

            result = run_script(SPATIAL, source, output)
            self.assertNotEqual(0, result.returncode)
            self.assertIn("SPATIAL_QA_PASS=false", result.stdout)
            self.assertIn("foreground-fill-occludes-content", result.stdout)

    def test_background_decoration_behind_text_passes_spatial_gate(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "good.pptx"
            make_source(source)

            prs = Presentation(source)
            slide = prs.slides[0]
            background = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(6), Inches(2))
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(240, 240, 240)
            background.line.fill.background()
            tree = slide.shapes._spTree
            tree.remove(background._element)
            tree.insert(2, background._element)
            prs.save(output)

            result = run_script(SPATIAL, source, output)
            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            self.assertIn("SPATIAL_QA_PASS=true", result.stdout)

    def test_composition_gate_requires_evidence_and_all_dimensions(self):
        with tempfile.TemporaryDirectory() as td:
            report_path = Path(td) / "composition.json"
            report = {
                "schema": 1,
                "slide_count": 1,
                "render_engine": "PowerPoint",
                "reviewer": "AI vision reviewer",
                "overall_pass": True,
                "slides": [
                    {
                        "slide": 1,
                        "composition_score": 92,
                        "checks": composition_checks(),
                        "scores": composition_scores(),
                        "evidence": {
                            "source_comparison": "brand and title compared side-by-side",
                            "grid_rails": ["left title rail", "right brand rail"],
                            "reading_order": ["title", "subtitle", "date"],
                            "brand_anchors": ["PEGATRON", "MEC"],
                        },
                    }
                ],
            }
            report_path.write_text(json.dumps(report), encoding="utf-8")
            good = run_script(COMPOSITION, report_path, "--expected-slides", "1")
            self.assertEqual(0, good.returncode, good.stdout + good.stderr)
            self.assertIn("COMPOSITION_QA_PASS=true", good.stdout)

            report["slides"][0]["checks"]["visual_balance_coherent"] = False
            report_path.write_text(json.dumps(report), encoding="utf-8")
            bad = run_script(COMPOSITION, report_path, "--expected-slides", "1")
            self.assertNotEqual(0, bad.returncode)
            self.assertIn("COMPOSITION_QA_PASS=false", bad.stdout)

    def test_composition_gate_rejects_beauty_score_that_hides_weak_dimension(self):
        with tempfile.TemporaryDirectory() as td:
            report_path = Path(td) / "composition.json"
            report = {
                "schema": 1,
                "slide_count": 1,
                "render_engine": "PowerPoint",
                "reviewer": "AI vision reviewer",
                "overall_pass": True,
                "slides": [
                    {
                        "slide": 1,
                        "composition_score": 94,
                        "checks": composition_checks(),
                        "scores": composition_scores(),
                        "evidence": {
                            "source_comparison": "source vs final",
                            "grid_rails": ["rail A"],
                            "reading_order": ["A", "B"],
                            "brand_anchors": ["brand"],
                        },
                    }
                ],
            }
            report["slides"][0]["scores"]["alignment"] = 70
            report_path.write_text(json.dumps(report), encoding="utf-8")
            result = run_script(COMPOSITION, report_path, "--expected-slides", "1")
            self.assertNotEqual(0, result.returncode)
            self.assertIn("alignment score 70", result.stdout)


if __name__ == "__main__":
    unittest.main()
