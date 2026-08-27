from __future__ import annotations

from pathlib import Path
import os
import sys
import tempfile
import unittest
from unittest import mock
import urllib.error

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
for path in (ROOT / "launcher", ROOT / "pptx-beautify-lock" / "scripts"):
    value = str(path)
    if value not in sys.path:
        sys.path.insert(0, value)

import offline_runtime
import update_manager


class OfflineOutputTests(unittest.TestCase):
    def test_success_means_final_file_exists_and_reopens(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            src = root / "source.pptx"
            out = root / "nested" / "final.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Final output regression"
            prs.save(src)

            logs: list[str] = []
            offline_runtime.beautify_to_final(
                src,
                out,
                "自動（忠於原稿 / Source-faithful）",
                logs.append,
                check_updates=False,
            )

            self.assertTrue(out.is_file())
            self.assertGreater(out.stat().st_size, 0)
            self.assertEqual(len(Presentation(out).slides), 1)
            self.assertIn("FINAL_OUTPUT_EXISTS=true", logs)
            self.assertIn("FINAL_OUTPUT_REOPEN_PASS=true", logs)
            self.assertEqual(logs.count("OFFLINE_BEAUTIFY_PASS=true"), 1)
            self.assertFalse(any(out.parent.glob("*.candidate.pptx")))


class OptionalUpdaterTests(unittest.TestCase):
    def test_no_network_skips_update_and_keeps_working(self):
        with tempfile.TemporaryDirectory() as td, mock.patch.dict(os.environ, {"LOCALAPPDATA": td}):
            with mock.patch.object(
                update_manager,
                "_request_json",
                side_effect=urllib.error.URLError("offline"),
            ):
                logs: list[str] = []
                selection = update_manager.check_for_engine_update(
                    local_version="0.7.1",
                    launcher_version="0.7.1",
                    log=logs.append,
                )
            self.assertIsNone(selection.engine_path)
            self.assertEqual(selection.engine_version, "0.7.1")
            self.assertEqual(selection.status, "offline_skip")
            self.assertIn("UPDATE_CHECK=offline_skip", logs)

    def test_newer_remote_engine_is_cached_and_reused_offline(self):
        fake_engine = b"STYLE_PRESETS = {}\ndef beautify_pptx(*args, **kwargs):\n    return None\n"
        branch_data = {"commit": {"sha": "abc123"}}
        manifest = {
            "engine_version": "0.7.2",
            "min_launcher_version": "0.7.1",
            "engine_path": "launcher/pptx_offline_engine.py",
        }
        with tempfile.TemporaryDirectory() as td, mock.patch.dict(os.environ, {"LOCALAPPDATA": td}):
            with mock.patch.object(update_manager, "_request_json", side_effect=[branch_data, manifest]), mock.patch.object(
                update_manager, "_request_bytes", return_value=fake_engine
            ):
                first = update_manager.check_for_engine_update(
                    local_version="0.7.1", launcher_version="0.7.1"
                )
            self.assertEqual(first.status, "updated")
            self.assertEqual(first.engine_version, "0.7.2")
            self.assertTrue(first.engine_path and first.engine_path.is_file())

            with mock.patch.object(
                update_manager,
                "_request_json",
                side_effect=urllib.error.URLError("offline"),
            ):
                second = update_manager.check_for_engine_update(
                    local_version="0.7.1", launcher_version="0.7.1"
                )
            self.assertEqual(second.engine_version, "0.7.2")
            self.assertTrue(second.engine_path and second.engine_path.is_file())

    def test_stale_update_channel_can_never_downgrade_engine(self):
        branch_data = {"commit": {"sha": "oldsha"}}
        manifest = {
            "engine_version": "0.6.2",
            "min_launcher_version": "0.6.2",
            "engine_path": "launcher/pptx_offline_engine.py",
        }
        with tempfile.TemporaryDirectory() as td, mock.patch.dict(os.environ, {"LOCALAPPDATA": td}):
            with mock.patch.object(update_manager, "_request_json", side_effect=[branch_data, manifest]):
                selection = update_manager.check_for_engine_update(
                    local_version="0.7.1", launcher_version="0.7.1"
                )
            self.assertEqual(selection.engine_version, "0.7.1")
            self.assertEqual(selection.status, "up_to_date")
            self.assertIsNone(selection.engine_path)


if __name__ == "__main__":
    unittest.main()
