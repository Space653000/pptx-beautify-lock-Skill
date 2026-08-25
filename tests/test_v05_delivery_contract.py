from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REGRESSION = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_regression.py"


def run_regression(*args: str):
    return subprocess.run(
        [sys.executable, str(REGRESSION), *map(str, args)],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


def make_clean_deck(path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(6.0), Inches(0.8))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "v0.5 delivery / 最終交付契約"
    run.font.size = Pt(24)
    prs.save(path)


def visual_report(path: Path):
    checks = {
        "no_unintended_overlap": True,
        "no_clipping_or_overflow": True,
        "content_visible": True,
        "text_readable": True,
        "hierarchy_clear": True,
        "alignment_consistent": True,
        "tables_charts_readable": True,
        "style_consistent": True,
        "no_template_placeholder_artifacts": True,
        "theme_fidelity_preserved": True,
        "bilingual_typography_clean": True,
    }
    payload = {
        "schema": 3,
        "slide_count": 1,
        "render_engine": "contract-test-renderer",
        "reviewer": "contract-test-reviewer",
        "overall_pass": True,
        "slides": [{"slide": 1, "score": 95, "checks": checks}],
    }
    path.write_text(json.dumps(payload), encoding="utf-8")


def composition_report(path: Path):
    checks = {
        "brand_chrome_respected": True,
        "content_not_occluded": True,
        "grid_alignment_coherent": True,
        "peer_components_aligned": True,
        "spacing_rhythm_coherent": True,
        "reading_order_clear": True,
        "visual_balance_coherent": True,
        "slide_role_composition_fit": True,
        "decorative_elements_earn_their_place": True,
    }
    scores = {
        "hierarchy": 93,
        "alignment": 94,
        "spacing": 92,
        "balance": 92,
        "brand_fidelity": 95,
        "restraint": 94,
        "data_legibility": 93,
    }
    payload = {
        "schema": 1,
        "slide_count": 1,
        "render_engine": "contract-test-renderer",
        "reviewer": "contract-test-reviewer",
        "overall_pass": True,
        "slides": [
            {
                "slide": 1,
                "composition_score": 93,
                "checks": checks,
                "scores": scores,
                "evidence": {
                    "source_comparison": "source and output compared side-by-side",
                    "grid_rails": ["title-left", "content-top"],
                    "reading_order": ["title"],
                    "brand_anchors": ["none in synthetic fixture; source canvas retained"],
                },
            }
        ],
    }
    path.write_text(json.dumps(payload), encoding="utf-8")


class V05DeliveryContractTests(unittest.TestCase):
    def test_v05_delivery_passes_only_with_both_render_reports(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "output.pptx"
            visual = Path(td) / "visual.json"
            composition = Path(td) / "composition.json"
            make_clean_deck(source)
            Presentation(source).save(output)
            visual_report(visual)
            composition_report(composition)

            result = run_regression(
                source,
                output,
                "--visual-qa-report",
                visual,
                "--require-visual-qa",
                "--composition-qa-report",
                composition,
                "--require-composition-qa",
            )
            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            self.assertIn("CONTENT_LOCK_PASS=true", result.stdout)
            self.assertIn("SPATIAL_QA_PASS=true", result.stdout)
            self.assertIn("VISUAL_QA_PASS=true", result.stdout)
            self.assertIn("COMPOSITION_QA_PASS=true", result.stdout)
            self.assertIn("REGRESSION_V05_PASS=true", result.stdout)
            self.assertIn("DELIVERY_V05_PASS=true", result.stdout)

    def test_v05_fails_closed_without_composition_report(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "output.pptx"
            visual = Path(td) / "visual.json"
            make_clean_deck(source)
            Presentation(source).save(output)
            visual_report(visual)

            result = run_regression(
                source,
                output,
                "--visual-qa-report",
                visual,
                "--require-visual-qa",
                "--require-composition-qa",
            )
            self.assertNotEqual(0, result.returncode)
            self.assertIn("COMPOSITION_QA_PASS=false", result.stdout)
            self.assertIn("REGRESSION_V05_PASS=false", result.stdout)
            self.assertIn("DELIVERY_V05_PASS=false", result.stdout)


if __name__ == "__main__":
    unittest.main()
