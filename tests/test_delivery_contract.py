from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REGRESSION = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_regression.py"


def run_regression(*args: str):
    return subprocess.run(
        [sys.executable, str(REGRESSION), *map(str, args)],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


def make_clean_deck(path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(6.0), Inches(0.8))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Delivery contract / 最終交付契約"
    run.font.size = Pt(24)
    prs.save(path)


def make_visual_report(path: Path):
    checks = {
        "no_unintended_overlap": True,
        "no_clipping_or_overflow": True,
        "content_visible": True,
        "text_readable": True,
        "hierarchy_clear": True,
        "alignment_consistent": True,
        "tables_charts_readable": True,
        "style_consistent": True,
        "no_template_placeholder_artifacts": True,
        "theme_fidelity_preserved": True,
        "bilingual_typography_clean": True,
    }
    payload = {
        "schema": 3,
        "slide_count": 1,
        "render_engine": "contract-test-renderer",
        "reviewer": "contract-test-reviewer",
        "overall_pass": True,
        "slides": [{"slide": 1, "score": 95, "checks": checks}],
    }
    path.write_text(json.dumps(payload), encoding="utf-8")


class DeliveryContractTests(unittest.TestCase):
    def test_clean_identical_deck_with_complete_visual_qa_passes_delivery(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "output.pptx"
            visual = Path(td) / "visual_qa.json"

            make_clean_deck(source)
            Presentation(source).save(output)
            make_visual_report(visual)

            result = run_regression(
                source,
                output,
                "--visual-qa-report",
                visual,
                "--require-visual-qa",
            )

            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            self.assertIn("CONTENT_LOCK_PASS=true", result.stdout)
            self.assertIn("THEME_GUARD_PASS=true", result.stdout)
            self.assertIn("THEME_FIDELITY_PASS=true", result.stdout)
            self.assertIn("LAYOUT_QA_PASS=true", result.stdout)
            self.assertIn("VISUAL_QA_PASS=true", result.stdout)
            self.assertIn("REGRESSION_PASS=true", result.stdout)
            self.assertIn("DELIVERY_PASS=true", result.stdout)


if __name__ == "__main__":
    unittest.main()
