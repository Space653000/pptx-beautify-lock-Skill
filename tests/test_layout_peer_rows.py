from __future__ import annotations

import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
SPATIAL = ROOT / "pptx-beautify-lock" / "scripts" / "pptx_layout_intelligence.py"


def run_spatial(source: Path, output: Path):
    return subprocess.run(
        [sys.executable, str(SPATIAL), str(source), str(output)],
        cwd=ROOT,
        text=True,
        capture_output=True,
        check=False,
    )


def chart_data():
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    data.add_series("S", (1, 2, 3))
    return data


class PeerRowRegressionTests(unittest.TestCase):
    def test_two_by_two_chart_grid_does_not_compare_across_rows(self):
        with tempfile.TemporaryDirectory() as td:
            source = Path(td) / "source.pptx"
            output = Path(td) / "output.pptx"

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            positions = [
                (0.6, 1.0),
                (6.9, 1.0),
                (0.6, 4.1),
                (6.9, 4.1),
            ]
            for left, top in positions:
                slide.shapes.add_chart(
                    XL_CHART_TYPE.LINE,
                    Inches(left),
                    Inches(top),
                    Inches(5.8),
                    Inches(2.5),
                    chart_data(),
                )
            prs.save(source)
            Presentation(source).save(output)

            result = run_spatial(source, output)
            self.assertEqual(0, result.returncode, result.stdout + result.stderr)
            self.assertIn("SPATIAL_QA_PASS=true", result.stdout)
            self.assertNotIn("sibling-visual-rail-drift", result.stdout)


if __name__ == "__main__":
    unittest.main()
