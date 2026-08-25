#!/usr/bin/env python3
"""Conservative PPTX geometry QA scanner.

Checks obvious out-of-bounds geometry, tiny text, and suspicious overlaps.
This is NOT a replacement for rendering and visual inspection.
"""

from __future__ import annotations

import argparse
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("LAYOUT_QA_PASS=false")
    print("ERROR=python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    raise SystemExit(3)


def box(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def area_intersection(a, b):
    x1 = max(a[0], b[0])
    y1 = max(a[1], b[1])
    x2 = min(a[2], b[2])
    y2 = min(a[3], b[3])
    if x2 <= x1 or y2 <= y1:
        return 0
    return (x2 - x1) * (y2 - y1)


def area(a):
    return max(0, a[2] - a[0]) * max(0, a[3] - a[1])


def has_visible_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return False
    return bool((shape.text or "").strip())


def min_font_pt(shape):
    values = []
    if not getattr(shape, "has_text_frame", False):
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                values.append(run.font.size.pt)
    return min(values) if values else None


def is_background_like(shape, sw, sh):
    # Full-slide or near-full-slide shapes are usually intentional backgrounds.
    a = box(shape)
    return area(a) >= 0.85 * sw * sh


def scan(path: str, overlap_threshold: float, tiny_pt: float):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    warnings = []
    hard_errors = []

    for sidx, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        for i, shape in enumerate(shapes):
            name = getattr(shape, "name", f"shape-{i}")
            if shape.width <= 0 or shape.height <= 0:
                hard_errors.append(f"slide {sidx}: non-positive geometry: {name}")
                continue
            l, t, r, b = box(shape)
            if l < 0 or t < 0 or r > sw or b > sh:
                hard_errors.append(f"slide {sidx}: out-of-bounds: {name}")
            m = min_font_pt(shape)
            if m is not None and m < tiny_pt:
                warnings.append(f"slide {sidx}: tiny text {m:.1f} pt: {name}")

        # Suspicious overlap heuristic: only flag pairs where at least one contains text,
        # ignore very large background-like shapes and groups.
        for i in range(len(shapes)):
            a_shape = shapes[i]
            if is_background_like(a_shape, sw, sh):
                continue
            if getattr(a_shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                continue
            for j in range(i + 1, len(shapes)):
                b_shape = shapes[j]
                if is_background_like(b_shape, sw, sh):
                    continue
                if getattr(b_shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                    continue
                if not (has_visible_text(a_shape) or has_visible_text(b_shape)):
                    continue
                ba, bb = box(a_shape), box(b_shape)
                inter = area_intersection(ba, bb)
                if inter <= 0:
                    continue
                denom = min(area(ba), area(bb))
                if denom <= 0:
                    continue
                ratio = inter / denom
                if ratio >= overlap_threshold:
                    warnings.append(
                        f"slide {sidx}: suspicious overlap {ratio:.1%}: "
                        f"{getattr(a_shape,'name','A')} <-> {getattr(b_shape,'name','B')}"
                    )

    return hard_errors, warnings, len(prs.slides)


def main():
    ap = argparse.ArgumentParser(description="PPTX layout QA / PPTX 版面品質掃描")
    ap.add_argument("pptx")
    ap.add_argument("--overlap-threshold", type=float, default=0.15,
                    help="minimum overlap ratio to warn (default: 0.15)")
    ap.add_argument("--tiny-pt", type=float, default=10.0,
                    help="font size below which to warn (default: 10pt)")
    args = ap.parse_args()

    try:
        hard, warnings, slides = scan(args.pptx, args.overlap_threshold, args.tiny_pt)
    except Exception as exc:
        print("LAYOUT_QA_PASS=false")
        print(f"ERROR={exc}", file=sys.stderr)
        return 3

    ok = not hard
    print(f"LAYOUT_QA_PASS={'true' if ok else 'false'}")
    print(f"slides_checked={slides}")
    print(f"layout_errors={len(hard)}")
    print(f"layout_warnings={len(warnings)}")
    for x in hard:
        print(f"ERROR: {x}")
    for x in warnings:
        print(f"WARNING: {x}")
    print("NOTE=Warnings require render/visual review; some overlaps are intentional.")
    return 0 if ok else 2


if __name__ == "__main__":
    raise SystemExit(main())
