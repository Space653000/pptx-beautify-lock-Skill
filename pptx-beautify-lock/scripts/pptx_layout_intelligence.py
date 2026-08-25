#!/usr/bin/env python3
"""Spatial composition guard for pptx-beautify-lock.

This module complements semantic Content Lock and render-based visual QA.
It detects a small set of high-confidence spatial regressions while leaving
subjective beauty decisions to rendered-slide review.

Key checks:
- foreground solid fills that occlude content;
- newly introduced large filled text regions over image-backed branded layouts;
- sibling chart/picture pairs whose tops/sizes drift on dense data slides;
- top-heavy data layouts that leave a large unused body region (review warning).
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import asdict, dataclass

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class Finding:
    slide: int
    severity: str
    rule: str
    message_zh_TW: str
    message_en: str
    objects: list[str]


def _box(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _area(box):
    return max(0, box[2] - box[0]) * max(0, box[3] - box[1])


def _intersection(a, b):
    x1, y1 = max(a[0], b[0]), max(a[1], b[1])
    x2, y2 = min(a[2], b[2]), min(a[3], b[3])
    return 0 if x2 <= x1 or y2 <= y1 else (x2 - x1) * (y2 - y1)


def _text(shape):
    if getattr(shape, "has_text_frame", False):
        return (shape.text or "").strip()
    if getattr(shape, "has_table", False):
        return "\n".join(
            (cell.text or "").strip()
            for row in shape.table.rows
            for cell in row.cells
        ).strip()
    return ""


def _is_content(shape):
    return (
        bool(_text(shape))
        or getattr(shape, "has_table", False)
        or getattr(shape, "shape_type", None)
        in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        }
    )


def _solid_fill(shape):
    try:
        return shape.fill.type == MSO_FILL.SOLID
    except Exception:
        return False


def _layout_full_bleed_image(slide, sw, sh):
    for scope in (slide.slide_layout.shapes, slide.slide_layout.slide_master.shapes):
        for obj in scope:
            try:
                if (
                    obj.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and _area(_box(obj)) >= 0.88 * sw * sh
                ):
                    return True
            except Exception:
                pass
    return False


def _role(slide):
    visuals = sum(
        1
        for shape in slide.shapes
        if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}
    )
    tables = sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))
    text_shapes = sum(1 for shape in slide.shapes if _text(shape))
    if tables or visuals >= 2:
        return "data"
    if text_shapes <= 2:
        return "cover_or_section"
    return "content"


def _signature(shape):
    text = _text(shape)
    if text:
        return ("text", text)
    if getattr(shape, "has_table", False):
        return ("table", text)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            return ("picture", shape.image.sha1)
        except Exception:
            return ("picture", shape.name)
    return ("kind", int(shape.shape_type), shape.name)


def _matched_source_map(slide):
    result = {}
    for shape in slide.shapes:
        result.setdefault(_signature(shape), []).append(shape)
    return result


def _result(errors, warnings):
    return {
        "schema": 1,
        "SPATIAL_QA_PASS": not any(item.severity == "ERROR" for item in errors),
        "errors": [asdict(item) for item in errors if item.severity == "ERROR"],
        "warnings": [asdict(item) for item in warnings],
        "error_count": sum(1 for item in errors if item.severity == "ERROR"),
        "warning_count": len(warnings),
        "RENDER_COMPOSITION_REVIEW_REQUIRED": bool(warnings),
    }


def compare_presentations(source, output):
    src = Presentation(source)
    out = Presentation(output)
    errors: list[Finding] = []
    warnings: list[Finding] = []

    if len(src.slides) != len(out.slides):
        errors.append(
            Finding(
                0,
                "ERROR",
                "slide-count-drift",
                "頁數已改變",
                "Slide count changed",
                [],
            )
        )
        return _result(errors, warnings)

    for idx, (source_slide, output_slide) in enumerate(zip(src.slides, out.slides), 1):
        sw, sh = out.slide_width, out.slide_height
        shapes = list(output_slide.shapes)

        # A. A later solid-filled object must not cover an earlier protected object.
        for foreground_index, foreground in enumerate(shapes):
            if not _solid_fill(foreground):
                continue
            if _area(_box(foreground)) < 0.005 * sw * sh:
                continue
            for background_index in range(foreground_index):
                background = shapes[background_index]
                if not _is_content(background):
                    continue
                intersection = _intersection(_box(background), _box(foreground))
                if intersection <= 0:
                    continue
                ratio = intersection / max(1, _area(_box(background)))
                threshold = 0.12 if bool(_text(background)) else 0.28
                if ratio >= threshold:
                    errors.append(
                        Finding(
                            idx,
                            "ERROR",
                            "foreground-fill-occludes-content",
                            f"前景實心物件遮住受保護內容約 {ratio:.0%}",
                            f"Foreground solid fill occludes protected content by about {ratio:.0%}",
                            [background.name, foreground.name],
                        )
                    )

        # B. A branded/full-bleed background plus a newly filled large content region
        # is not automatically wrong, but it is always a render-comparison risk.
        if _layout_full_bleed_image(output_slide, sw, sh):
            source_map = _matched_source_map(source_slide)
            for output_shape in output_slide.shapes:
                if not _solid_fill(output_shape):
                    continue
                area_ratio = _area(_box(output_shape)) / (sw * sh)
                if area_ratio < 0.10:
                    continue
                source_matches = source_map.get(_signature(output_shape), [])
                if source_matches and all(not _solid_fill(item) for item in source_matches):
                    warnings.append(
                        Finding(
                            idx,
                            "WARNING",
                            "brand-background-occlusion-risk",
                            f"品牌/滿版圖片背景上新增大面積實心內容區塊（{area_ratio:.0%}），必須以 source-vs-final render 確認沒有壓住品牌識別或關鍵背景",
                            f"Large solid content region ({area_ratio:.0%}) was added over a branded/full-bleed image layout; source-vs-final render must confirm brand anchors remain clear",
                            [output_shape.name],
                        )
                    )

        # C. On technical/data slides, peer visuals should share rails and size.
        if _role(output_slide) == "data":
            visuals = [
                shape
                for shape in output_slide.shapes
                if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}
                and _area(_box(shape)) >= 0.05 * sw * sh
            ]
            for first_index in range(len(visuals)):
                for second_index in range(first_index + 1, len(visuals)):
                    first = visuals[first_index]
                    second = visuals[second_index]
                    first_area, second_area = _area(_box(first)), _area(_box(second))
                    similarity = min(first_area, second_area) / max(first_area, second_area)
                    if similarity < 0.72:
                        continue
                    first_center_x = (first.left + first.width / 2) / sw
                    second_center_x = (second.left + second.width / 2) / sw
                    if abs(first_center_x - second_center_x) < 0.18:
                        continue
                    top_drift = abs(first.top - second.top) / sh
                    height_drift = abs(first.height - second.height) / sh
                    width_drift = abs(first.width - second.width) / sw
                    if top_drift > 0.025 or height_drift > 0.045 or width_drift > 0.045:
                        errors.append(
                            Finding(
                                idx,
                                "ERROR",
                                "sibling-visual-rail-drift",
                                "同列、同角色的大型圖表/圖片沒有共用上緣與尺寸骨架",
                                "Peer visuals in the same row do not share a coherent top rail and size",
                                [first.name, second.name],
                            )
                        )

            # D. A dense data slide that ends too early is a balance review risk.
            majors = [
                shape
                for shape in output_slide.shapes
                if _is_content(shape) and _area(_box(shape)) >= 0.02 * sw * sh
            ]
            if majors:
                bottom = max((shape.top + shape.height) / sh for shape in majors)
                if bottom < 0.79 and len(majors) >= 4:
                    warnings.append(
                        Finding(
                            idx,
                            "WARNING",
                            "data-body-vertical-balance-risk",
                            f"主要資料內容在頁高 {bottom:.0%} 前即結束，底部留白偏大；應檢查是否過度擠在上半部",
                            f"Primary data content ends by {bottom:.0%} of slide height; review whether the body is unnecessarily top-heavy",
                            [shape.name for shape in majors[:6]],
                        )
                    )

    return _result(errors, warnings)


def main() -> int:
    parser = argparse.ArgumentParser(description="PPTX Spatial QA / PPTX 空間骨骼檢查")
    parser.add_argument("source")
    parser.add_argument("output")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args()

    try:
        result = compare_presentations(args.source, args.output)
    except Exception as exc:
        print("SPATIAL_QA_PASS=false")
        print(f"ERROR={exc}", file=sys.stderr)
        return 3

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print(f"SPATIAL_QA_PASS={'true' if result['SPATIAL_QA_PASS'] else 'false'}")
        print(f"spatial_errors={result['error_count']}")
        print(f"spatial_warnings={result['warning_count']}")
        print(
            "RENDER_COMPOSITION_REVIEW_REQUIRED="
            + ("true" if result["RENDER_COMPOSITION_REVIEW_REQUIRED"] else "false")
        )
        for item in result["errors"] + result["warnings"]:
            print(
                f"{item['severity']}: slide={item['slide']} rule={item['rule']} "
                f"{item['message_zh_TW']}"
            )
    return 0 if result["SPATIAL_QA_PASS"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
