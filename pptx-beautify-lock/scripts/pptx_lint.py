#!/usr/bin/env python3
"""PPTX Linter for pptx-beautify-lock.

Scans geometry, typography, table density, overlap, template-placeholder
leakage, CJK fallback risks, edge margins, and cross-slide consistency.

This tool never modifies the input file. Rendered Visual QA remains the final
authority for appearance, source-theme fidelity, and actual font rendering.
"""

from __future__ import annotations

import argparse
from collections import Counter
import json
import statistics
import sys
from dataclasses import asdict, dataclass

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR=python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    raise SystemExit(3)

EMU_PER_INCH = 914400

GENERIC_TEMPLATE_TEXTS = {
    "presentation title",
    "presentation subtitle",
    "click to add title",
    "click to add subtitle",
    "click to add text",
    "click to add content",
    "title placeholder",
    "subtitle placeholder",
}

# These fonts are fine for pure Latin text, but should not be trusted as the
# sole explicit family for Traditional-Chinese/mixed runs. PowerPoint may fall
# back silently to another CJK font and break visual consistency.
LATIN_ONLY_RISK_FONTS = {
    "arial",
    "helvetica",
    "inter",
    "aptos",
    "aptos display",
    "calibri",
    "roboto",
}


@dataclass
class Finding:
    slide: int
    severity: str
    rule: str
    message_zh_TW: str
    message_en: str
    objects: list[str]


def _box(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _area(b):
    return max(0, b[2] - b[0]) * max(0, b[3] - b[1])


def _area_in2(shape):
    return max(0.0, shape.width / EMU_PER_INCH) * max(0.0, shape.height / EMU_PER_INCH)


def _intersection(a, b):
    x1, y1 = max(a[0], b[0]), max(a[1], b[1])
    x2, y2 = min(a[2], b[2]), min(a[3], b[3])
    if x2 <= x1 or y2 <= y1:
        return 0
    return (x2 - x1) * (y2 - y1)


def _visible_text(shape):
    if getattr(shape, "has_text_frame", False) and (shape.text or "").strip():
        return True
    if getattr(shape, "has_table", False):
        return any((cell.text or "").strip() for row in shape.table.rows for cell in row.cells)
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return any(_visible_text(child) for child in shape.shapes)
    return False


def _normalize_text(value: str) -> str:
    return " ".join((value or "").strip().casefold().split())


def _generic_template_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    text = _normalize_text(getattr(shape, "text", ""))
    return text if text in GENERIC_TEMPLATE_TEXTS else None


def _contains_cjk(text: str) -> bool:
    for ch in text or "":
        code = ord(ch)
        if (
            0x3400 <= code <= 0x4DBF
            or 0x4E00 <= code <= 0x9FFF
            or 0xF900 <= code <= 0xFAFF
            or 0x3100 <= code <= 0x312F
        ):
            return True
    return False


def _iter_text_frames(shape, prefix: str | None = None):
    name = prefix or getattr(shape, "name", "shape")
    if getattr(shape, "has_text_frame", False):
        yield name, shape.text_frame

    if getattr(shape, "has_table", False):
        for r_idx, row in enumerate(shape.table.rows, 1):
            for c_idx, cell in enumerate(row.cells, 1):
                yield f"{name}[R{r_idx}C{c_idx}]", cell.text_frame

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            child_name = f"{name}/{getattr(child, 'name', 'child')}"
            yield from _iter_text_frames(child, child_name)


def _font_sizes(shape):
    values = []
    for _, tf in _iter_text_frames(shape):
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    values.append(float(r.font.size.pt))
    return values


def _font_names(shape):
    names = set()
    for _, tf in _iter_text_frames(shape):
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.name:
                    names.add(r.font.name.strip())
    return names


def _font_fallback_risks(shape):
    risks = []
    for location, tf in _iter_text_frames(shape):
        for p_idx, paragraph in enumerate(tf.paragraphs, 1):
            for r_idx, run in enumerate(paragraph.runs, 1):
                text = run.text or ""
                font_name = (run.font.name or "").strip()
                if not font_name or not _contains_cjk(text):
                    continue
                if font_name.casefold() in LATIN_ONLY_RISK_FONTS:
                    risks.append((location, p_idx, r_idx, font_name, text[:40]))
    return risks


def _text_chars(shape):
    total = 0
    for _, tf in _iter_text_frames(shape):
        total += sum(len((p.text or "").strip()) for p in tf.paragraphs)
    return total


def _background_like(shape, sw, sh):
    return _area(_box(shape)) >= 0.85 * sw * sh


def _table_density(shape):
    if not getattr(shape, "has_table", False):
        return None
    table = shape.table
    rows = len(table.rows)
    cols = len(table.columns)
    cells = rows * cols
    if cells <= 0:
        return None
    chars = sum(len((cell.text or "").strip()) for row in table.rows for cell in row.cells)
    area_in2 = max(_area_in2(shape), 0.01)
    return {
        "rows": rows,
        "cols": cols,
        "cells": cells,
        "chars": chars,
        "avg_cell_area_in2": area_in2 / cells,
        "chars_per_in2": chars / area_in2,
    }


def _title_profile(shape, sw, sh):
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        ph_name = getattr(shape.placeholder_format.type, "name", str(shape.placeholder_format.type))
    except Exception:
        return None
    if ph_name != "TITLE":
        return None
    return (
        shape.left / sw,
        shape.top / sh,
        shape.width / sw,
        shape.height / sh,
        getattr(shape, "name", "title"),
    )


def scan_presentation(path: str, tiny_pt: float = 11.0, overlap_threshold: float = 0.15,
                      edge_margin_ratio: float = 0.01, max_fonts_per_slide: int = 4):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    findings: list[Finding] = []
    slide_font_sets: dict[int, set[str]] = {}
    title_profiles: list[tuple[int, tuple]] = []

    for sidx, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        slide_fonts = set()
        template_indexes: set[int] = set()

        for i, shape in enumerate(shapes):
            name = getattr(shape, "name", f"shape-{i}")
            l, t, r, b = _box(shape)

            if shape.width <= 0 or shape.height <= 0:
                findings.append(Finding(sidx, "ERROR", "non-positive-geometry",
                    "物件寬度或高度不是正值", "Object has non-positive geometry", [name]))
                continue

            if l < 0 or t < 0 or r > sw or b > sh:
                findings.append(Finding(sidx, "ERROR", "out-of-bounds",
                    "物件超出投影片邊界", "Object extends beyond slide bounds", [name]))

            if not _background_like(shape, sw, sh):
                mx, my = sw * edge_margin_ratio, sh * edge_margin_ratio
                if l < mx or t < my or (sw - r) < mx or (sh - b) < my:
                    findings.append(Finding(sidx, "INFO", "unsafe-edge-margin",
                        "物件非常接近投影片邊界，請確認是否為刻意設計",
                        "Object is very close to a slide edge; confirm this is intentional", [name]))

            generic = _generic_template_text(shape)
            if generic:
                template_indexes.add(i)
                if getattr(shape, "is_placeholder", False):
                    findings.append(Finding(sidx, "ERROR", "template-placeholder-artifact",
                        f"偵測到 final 不應顯示的模板 placeholder 文字：{generic!r}",
                        f"Generic template placeholder text is visible: {generic!r}", [name]))
                else:
                    findings.append(Finding(sidx, "WARNING", "generic-template-text",
                        f"偵測到疑似模板示意文字：{generic!r}，需確認是否真為使用者內容",
                        f"Generic template-like text detected: {generic!r}; confirm whether it is user-authored",
                        [name]))

            sizes = _font_sizes(shape)
            if sizes and min(sizes) < tiny_pt:
                findings.append(Finding(sidx, "WARNING", "tiny-text",
                    f"偵測到 {min(sizes):.1f} pt 的小字，可能影響投影可讀性",
                    f"Detected {min(sizes):.1f} pt text; projected readability may be poor", [name]))

            for location, p_idx, r_idx, font_name, preview in _font_fallback_risks(shape):
                findings.append(Finding(sidx, "WARNING", "cjk-font-fallback-risk",
                    f"繁中/中文 run 明確指定為偏 Latin 字體 {font_name!r}，可能產生不可控 fallback",
                    f"CJK text is explicitly assigned to Latin-oriented font {font_name!r}; fallback may be unstable",
                    [f"{location}/P{p_idx}/R{r_idx}: {preview}"]))

            slide_fonts |= _font_names(shape)

            density = _table_density(shape)
            if density is not None:
                high_density = (
                    (density["cells"] >= 8 and density["avg_cell_area_in2"] < 0.22)
                    or density["chars_per_in2"] > 85
                )
                if high_density:
                    findings.append(Finding(sidx, "WARNING", "table-density-risk",
                        "表格資訊密度偏高，請檢查欄寬、列高、padding 與實際 render 可讀性",
                        "Table density is high; review sizing, padding and rendered readability",
                        [name]))

            chars = _text_chars(shape)
            area_in2 = _area_in2(shape)
            if chars >= 280 and area_in2 > 0 and chars / area_in2 > 55:
                findings.append(Finding(sidx, "INFO", "dense-text-region",
                    "文字區域資訊密度偏高，需 render 確認沒有 overflow/clipping",
                    "Dense text region; render to confirm there is no overflow/clipping", [name]))

            title = _title_profile(shape, sw, sh)
            if title is not None:
                title_profiles.append((sidx, title))

        slide_font_sets[sidx] = slide_fonts

        if len(slide_fonts) > max_fonts_per_slide:
            findings.append(Finding(sidx, "WARNING", "too-many-fonts",
                f"同頁偵測到 {len(slide_fonts)} 種明確字型，視覺一致性風險偏高",
                f"Detected {len(slide_fonts)} explicit font families on one slide", sorted(slide_fonts)))

        for i in range(len(shapes)):
            a = shapes[i]
            if _background_like(a, sw, sh) or getattr(a, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                continue
            for j in range(i + 1, len(shapes)):
                b = shapes[j]
                if _background_like(b, sw, sh) or getattr(b, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                    continue
                if not (_visible_text(a) or _visible_text(b)):
                    continue
                ba, bb = _box(a), _box(b)
                inter = _intersection(ba, bb)
                if inter <= 0:
                    continue
                denom = min(_area(ba), _area(bb))
                if denom <= 0:
                    continue
                ratio = inter / denom

                if (i in template_indexes or j in template_indexes) and _visible_text(a) and _visible_text(b):
                    findings.append(Finding(sidx, "ERROR", "template-artifact-overlap",
                        f"模板示意文字與真正可見內容重疊 {ratio:.1%}；必須保留真正內容並停用模板 artifact",
                        f"Generic template text overlaps visible content by {ratio:.1%}; preserve real content and suppress the artifact",
                        [getattr(a, "name", "A"), getattr(b, "name", "B")]))
                    continue

                if ratio >= overlap_threshold:
                    findings.append(Finding(sidx, "WARNING", "suspicious-overlap",
                        f"兩個物件疑似重疊 {ratio:.1%}，需 render 確認是否為刻意",
                        f"Two objects overlap by {ratio:.1%}; render review is required",
                        [getattr(a, "name", "A"), getattr(b, "name", "B")]))

    if len(prs.slides) >= 3:
        font_counts = Counter(font for fonts in slide_font_sets.values() for font in fonts)
        if len(font_counts) > 2:
            for sidx, fonts in slide_font_sets.items():
                outliers = sorted(font for font in fonts if font_counts[font] == 1)
                if outliers:
                    findings.append(Finding(sidx, "INFO", "cross-slide-font-outlier",
                        "此頁含其他頁未使用的明確字型，請確認是否為刻意",
                        "This slide uses explicit font families not used on other slides; confirm intentionally",
                        outliers))

    if len(title_profiles) >= 3:
        med_left = statistics.median(p[1][0] for p in title_profiles)
        med_top = statistics.median(p[1][1] for p in title_profiles)
        med_width = statistics.median(p[1][2] for p in title_profiles)
        med_height = statistics.median(p[1][3] for p in title_profiles)
        for sidx, profile in title_profiles:
            left, top, width, height, name = profile
            if (
                abs(left - med_left) > 0.04
                or abs(top - med_top) > 0.04
                or abs(width - med_width) > 0.10
                or abs(height - med_height) > 0.08
            ):
                findings.append(Finding(sidx, "INFO", "title-layout-outlier",
                    "標題區位置/尺寸與多數一般頁差異較大，請確認跨頁設計一致性",
                    "Title geometry differs materially from most standard slides; review cross-slide consistency",
                    [name]))

    errors = sum(1 for f in findings if f.severity == "ERROR")
    warnings = sum(1 for f in findings if f.severity == "WARNING")
    infos = sum(1 for f in findings if f.severity == "INFO")
    by_rule = Counter(f.rule for f in findings)
    return {
        "slides_checked": len(prs.slides),
        "lint_errors": errors,
        "lint_warnings": warnings,
        "lint_infos": infos,
        "findings_by_rule": dict(sorted(by_rule.items())),
        "LINT_PASS": errors == 0,
        "findings": [asdict(f) for f in findings],
    }


def main() -> int:
    ap = argparse.ArgumentParser(description="PPTX Linter / PowerPoint 版面檢查器")
    ap.add_argument("pptx")
    ap.add_argument("--tiny-pt", type=float, default=11.0)
    ap.add_argument("--overlap-threshold", type=float, default=0.15)
    ap.add_argument("--edge-margin-ratio", type=float, default=0.01)
    ap.add_argument("--max-fonts-per-slide", type=int, default=4)
    ap.add_argument("--json", action="store_true", help="Output full JSON findings")
    args = ap.parse_args()

    try:
        result = scan_presentation(
            args.pptx,
            tiny_pt=args.tiny_pt,
            overlap_threshold=args.overlap_threshold,
            edge_margin_ratio=args.edge_margin_ratio,
            max_fonts_per_slide=args.max_fonts_per_slide,
        )
    except Exception as exc:
        print("LINT_PASS=false")
        print(f"ERROR={exc}", file=sys.stderr)
        return 3

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print(f"LINT_PASS={'true' if result['LINT_PASS'] else 'false'}")
        print(f"slides_checked={result['slides_checked']}")
        print(f"lint_errors={result['lint_errors']}")
        print(f"lint_warnings={result['lint_warnings']}")
        print(f"lint_infos={result['lint_infos']}")
        for f in result["findings"]:
            objs = ", ".join(f["objects"])
            print(f"{f['severity']}: slide {f['slide']} [{f['rule']}] {f['message_zh_TW']} ({objs})")

    return 0 if result["LINT_PASS"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
