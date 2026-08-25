#!/usr/bin/env python3
"""Discover and compare source visual DNA for pptx-beautify-lock.

繁體中文：在美化前辨識來源簡報的 light/dark/mixed canvas、主題色、
大面積色塊與來源字型；美化後攔截高信心的色系極性翻轉。

English: Profile a source deck's canvas polarity, theme colors, large-area
visual mass, and fonts; compare output against source to block high-confidence
source-theme inversions.

This is a conservative structural heuristic. Render Visual QA remains the
final authority for source-theme fidelity.
"""

from __future__ import annotations

import argparse
from collections import Counter
import json
import sys
import zipfile
import xml.etree.ElementTree as ET

try:
    from pptx import Presentation
    from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
except ImportError:
    print("ERROR=python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    raise SystemExit(3)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"a": A_NS}
PROFILE_SCHEMA = 1

THEME_ENUM_TO_KEY = {
    "ACCENT_1": "accent1",
    "ACCENT_2": "accent2",
    "ACCENT_3": "accent3",
    "ACCENT_4": "accent4",
    "ACCENT_5": "accent5",
    "ACCENT_6": "accent6",
    "DARK_1": "dk1",
    "DARK_2": "dk2",
    "LIGHT_1": "lt1",
    "LIGHT_2": "lt2",
    "TEXT_1": "dk1",
    "TEXT_2": "dk2",
    "BACKGROUND_1": "lt1",
    "BACKGROUND_2": "lt2",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
}


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _normalize_hex(value: str | None) -> str | None:
    if value is None:
        return None
    value = str(value).strip().lstrip("#").upper()
    if len(value) == 6 and all(ch in "0123456789ABCDEF" for ch in value):
        return value
    return None


def _rgb_tuple(hex_value: str):
    return tuple(int(hex_value[i:i + 2], 16) for i in (0, 2, 4))


def _luminance(hex_value: str) -> float:
    r, g, b = _rgb_tuple(hex_value)
    return (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255.0


def _saturation(hex_value: str) -> float:
    r, g, b = (x / 255.0 for x in _rgb_tuple(hex_value))
    mx, mn = max(r, g, b), min(r, g, b)
    return 0.0 if mx == mn else (mx - mn) / (1.0 - abs(mx + mn - 1.0) or 1.0)


def _mode_for_color(hex_value: str | None) -> str:
    if not hex_value:
        return "unknown"
    lum = _luminance(hex_value)
    if lum >= 0.70:
        return "light"
    if lum <= 0.30:
        return "dark"
    return "mid"


def _read_theme(path: str):
    colors: dict[str, str] = {}
    fonts: dict[str, str] = {}
    try:
        with zipfile.ZipFile(path, "r") as zf:
            theme_names = [n for n in zf.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
            if not theme_names:
                return colors, fonts
            root = ET.fromstring(zf.read(sorted(theme_names)[0]))
    except Exception:
        return colors, fonts

    scheme = root.find(".//a:clrScheme", NS)
    if scheme is not None:
        for item in list(scheme):
            key = _local_name(item.tag)
            if not list(item):
                continue
            color_node = list(item)[0]
            kind = _local_name(color_node.tag)
            value = None
            if kind == "srgbClr":
                value = color_node.attrib.get("val")
            elif kind == "sysClr":
                value = color_node.attrib.get("lastClr") or color_node.attrib.get("val")
            norm = _normalize_hex(value)
            if norm:
                colors[key] = norm

    font_scheme = root.find(".//a:fontScheme", NS)
    if font_scheme is not None:
        for role in ("majorFont", "minorFont"):
            node = font_scheme.find(f"./a:{role}", NS)
            if node is None:
                continue
            for script in ("latin", "ea", "cs"):
                fnode = node.find(f"./a:{script}", NS)
                if fnode is not None and fnode.attrib.get("typeface"):
                    fonts[f"{role}.{script}"] = fnode.attrib["typeface"]
    return colors, fonts


def _color_from_format(color_format, theme_colors: dict[str, str]) -> str | None:
    try:
        ctype = color_format.type
    except Exception:
        return None
    try:
        if ctype == MSO_COLOR_TYPE.RGB:
            return _normalize_hex(str(color_format.rgb))
    except Exception:
        pass
    try:
        if ctype == MSO_COLOR_TYPE.SCHEME:
            enum_name = getattr(color_format.theme_color, "name", str(color_format.theme_color))
            key = THEME_ENUM_TO_KEY.get(enum_name)
            if key:
                return theme_colors.get(key)
    except Exception:
        pass
    return None


def _solid_fill_rgb(fill, theme_colors: dict[str, str]) -> str | None:
    try:
        if fill.type != MSO_FILL.SOLID:
            return None
        return _color_from_format(fill.fore_color, theme_colors)
    except Exception:
        return None


def _background_candidates(slide, theme_colors: dict[str, str]):
    candidates = []
    for label, holder, confidence in (
        ("slide-background", slide, 0.92),
        ("layout-background", getattr(slide, "slide_layout", None), 0.82),
        ("master-background", getattr(getattr(slide, "slide_layout", None), "slide_master", None), 0.76),
    ):
        if holder is None:
            continue
        try:
            rgb = _solid_fill_rgb(holder.background.fill, theme_colors)
        except Exception:
            rgb = None
        if rgb:
            candidates.append((confidence, label, rgb))
    return candidates


def _shape_area_ratio(shape, sw, sh) -> float:
    try:
        if shape.width <= 0 or shape.height <= 0:
            return 0.0
        return min(1.0, (shape.width * shape.height) / float(sw * sh))
    except Exception:
        return 0.0


def _shape_fill_rgb(shape, theme_colors: dict[str, str]) -> str | None:
    try:
        return _solid_fill_rgb(shape.fill, theme_colors)
    except Exception:
        return None


def _font_counter(prs: Presentation) -> Counter:
    counter = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            frames = []
            if getattr(shape, "has_text_frame", False):
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        frames.append(cell.text_frame)
            for tf in frames:
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            counter[run.font.name.strip()] += max(1, len(run.text or ""))
    return counter


def _slide_profile(slide, theme_colors: dict[str, str], sw, sh):
    candidates = _background_candidates(slide, theme_colors)
    dark_mass = 0.0
    colored_mass = 0.0
    fill_counter = Counter()

    for shape in slide.shapes:
        area = _shape_area_ratio(shape, sw, sh)
        if area <= 0:
            continue
        rgb = _shape_fill_rgb(shape, theme_colors)
        if not rgb:
            continue
        fill_counter[rgb] += round(area, 4)
        lum = _luminance(rgb)
        sat = _saturation(rgb)
        if lum <= 0.30:
            dark_mass += area
        if sat >= 0.18 and 0.12 < lum < 0.92:
            colored_mass += area
        if area >= 0.80:
            # A nearly full-page solid shape is stronger evidence than inherited
            # master defaults because it is what the viewer actually sees.
            candidates.append((0.97, "large-area-fill", rgb))
        elif area >= 0.55:
            candidates.append((0.88, "dominant-area-fill", rgb))

    if candidates:
        confidence, source, background = max(candidates, key=lambda item: item[0])
    else:
        # Standard PowerPoint decks with no explicit background render on the
        # light canvas. Confidence is deliberately moderate; render QA must
        # confirm custom image/master cases.
        confidence, source, background = 0.56, "default-light-fallback", "FFFFFF"

    mode = _mode_for_color(background)
    return {
        "canvas_mode": mode,
        "confidence": round(confidence, 3),
        "background_rgb": background,
        "background_evidence": source,
        "dark_fill_area_ratio": round(min(1.0, dark_mass), 3),
        "colored_fill_area_ratio": round(min(1.0, colored_mass), 3),
        "dominant_fill_colors": [
            {"rgb": rgb, "area_weight": weight}
            for rgb, weight in fill_counter.most_common(6)
        ],
    }


def profile_presentation(path: str) -> dict:
    prs = Presentation(path)
    theme_colors, theme_fonts = _read_theme(path)
    slides = [
        {"slide": idx, **_slide_profile(slide, theme_colors, prs.slide_width, prs.slide_height)}
        for idx, slide in enumerate(prs.slides, 1)
    ]

    decisive = [s for s in slides if s["canvas_mode"] in {"light", "dark"}]
    light = sum(1 for s in decisive if s["canvas_mode"] == "light")
    dark = sum(1 for s in decisive if s["canvas_mode"] == "dark")
    total = max(1, len(decisive))
    if light / total >= 0.75:
        deck_mode = "light"
    elif dark / total >= 0.75:
        deck_mode = "dark"
    elif light and dark:
        deck_mode = "mixed"
    else:
        deck_mode = "unknown"

    explicit_fonts = _font_counter(prs)
    accent_candidates = []
    for key in ("accent1", "accent2", "accent3", "accent4", "accent5", "accent6"):
        if key in theme_colors:
            accent_candidates.append({"source": key, "rgb": theme_colors[key]})

    avg_conf = sum(s["confidence"] for s in slides) / max(1, len(slides))
    review_required = deck_mode in {"mixed", "unknown"} or any(s["confidence"] < 0.70 for s in slides)

    return {
        "schema": PROFILE_SCHEMA,
        "canvas_mode": deck_mode,
        "confidence": round(avg_conf, 3),
        "slides": slides,
        "theme_colors": theme_colors,
        "theme_fonts": theme_fonts,
        "source_fonts": [
            {"name": name, "weight": weight}
            for name, weight in explicit_fonts.most_common(10)
        ],
        "accent_candidates": accent_candidates,
        "review_required": review_required,
    }


def compare_profiles(source_profile: dict, output_profile: dict) -> dict:
    violations = []
    warnings = []
    source_slides = {item["slide"]: item for item in source_profile.get("slides", [])}
    output_slides = {item["slide"]: item for item in output_profile.get("slides", [])}

    for slide_no in sorted(set(source_slides) & set(output_slides)):
        src = source_slides[slide_no]
        out = output_slides[slide_no]
        src_mode = src.get("canvas_mode")
        out_mode = out.get("canvas_mode")
        src_conf = float(src.get("confidence", 0))
        out_conf = float(out.get("confidence", 0))

        if (
            src_mode in {"light", "dark"}
            and out_mode in {"light", "dark"}
            and src_mode != out_mode
            and src_conf >= 0.55
            and out_conf >= 0.70
        ):
            violations.append({
                "slide": slide_no,
                "rule": "canvas-polarity-inversion",
                "source": src_mode,
                "output": out_mode,
                "message_zh_TW": "高信心偵測到來源頁面的明暗色系被翻轉",
                "message_en": "High-confidence source canvas polarity inversion detected",
            })

        src_dark = float(src.get("dark_fill_area_ratio", 0))
        out_dark = float(out.get("dark_fill_area_ratio", 0))
        if src_mode == "light" and src_dark <= 0.20 and out_dark - src_dark >= 0.38:
            violations.append({
                "slide": slide_no,
                "rule": "dark-visual-mass-drift",
                "source_dark_area_ratio": src_dark,
                "output_dark_area_ratio": out_dark,
                "message_zh_TW": "來源為淺色頁面，但輸出新增過多大面積深色視覺量",
                "message_en": "Output adds excessive dark visual mass to a light source slide",
            })

        if src_conf < 0.70 or out_conf < 0.70:
            warnings.append({
                "slide": slide_no,
                "rule": "render-theme-review-required",
                "message_zh_TW": "此頁結構式主色判定信心不足，必須以 render 視覺確認",
                "message_en": "Structural theme confidence is low; rendered review is required",
            })

    source_mode = source_profile.get("canvas_mode")
    output_mode = output_profile.get("canvas_mode")
    if (
        source_mode in {"light", "dark"}
        and output_mode in {"light", "dark"}
        and source_mode != output_mode
        and float(source_profile.get("confidence", 0)) >= 0.65
        and float(output_profile.get("confidence", 0)) >= 0.65
    ):
        violations.append({
            "slide": None,
            "rule": "deck-canvas-polarity-inversion",
            "source": source_mode,
            "output": output_mode,
            "message_zh_TW": "整份簡報主明暗色系被翻轉",
            "message_en": "Deck-level canvas polarity was inverted",
        })

    return {
        "THEME_GUARD_PASS": len(violations) == 0,
        "THEME_REVIEW_REQUIRED": bool(warnings) or source_profile.get("review_required", True) or output_profile.get("review_required", True),
        "violations": violations,
        "warnings": warnings,
        "source_canvas_mode": source_mode,
        "output_canvas_mode": output_mode,
    }


def compare_presentations(source: str, output: str) -> dict:
    return compare_profiles(profile_presentation(source), profile_presentation(output))


def main() -> int:
    ap = argparse.ArgumentParser(description="PPTX Theme Discovery / 來源主色調辨識與守門")
    sub = ap.add_subparsers(dest="command", required=True)

    p_profile = sub.add_parser("profile", help="Build source theme profile")
    p_profile.add_argument("pptx")
    p_profile.add_argument("--out")
    p_profile.add_argument("--json", action="store_true")

    p_compare = sub.add_parser("compare", help="Compare source/output theme polarity")
    p_compare.add_argument("source")
    p_compare.add_argument("output")
    p_compare.add_argument("--json", action="store_true")

    args = ap.parse_args()

    try:
        if args.command == "profile":
            result = profile_presentation(args.pptx)
            if args.out:
                with open(args.out, "w", encoding="utf-8") as f:
                    json.dump(result, f, ensure_ascii=False, indent=2)
            if args.json or not args.out:
                print(json.dumps(result, ensure_ascii=False, indent=2))
            else:
                print(f"THEME_PROFILE_WRITTEN={args.out}")
                print(f"canvas_mode={result['canvas_mode']}")
                print(f"theme_review_required={'true' if result['review_required'] else 'false'}")
            return 0

        result = compare_presentations(args.source, args.output)
        if args.json:
            print(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            print(f"THEME_GUARD_PASS={'true' if result['THEME_GUARD_PASS'] else 'false'}")
            print(f"THEME_REVIEW_REQUIRED={'true' if result['THEME_REVIEW_REQUIRED'] else 'false'}")
            print(f"theme_violations={len(result['violations'])}")
            for item in result["violations"]:
                print(f"ERROR: slide={item.get('slide')} rule={item['rule']} {item['message_zh_TW']}")
        return 0 if result["THEME_GUARD_PASS"] else 2
    except Exception as exc:
        print("THEME_GUARD_PASS=false")
        print(f"ERROR={exc}", file=sys.stderr)
        return 3


if __name__ == "__main__":
    raise SystemExit(main())
